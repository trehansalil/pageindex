"""Leaf format converters: LibreOffice, HTML, XLSX, image, rasterize, OCR, VLM, DOCX, PPTX.

Mechanical extraction from converters.py — lines 3769-4213 plus 3997-4009.
"""

from __future__ import annotations

import asyncio
import logging
import os
import re
import shutil
import subprocess
import tempfile

from ..script import normalize_dashes
from .docling_conv import _docling_converter, _repair_docling_tables

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants (lines 3997-4001)
# ---------------------------------------------------------------------------

# D4 (RFC-024): fallback rasterization backend for tesseract_ocr_pdf_pages. CMap
# corruption that crashes pypdfium2 page rendering is deterministic per-page, so a
# retry against pypdfium2 would fail identically; fitz uses a different rendering
# path (already proven for crop rasterization in _recover_picture_text) and isolates
# D7's rasterization from rasterize_pdf_pages' shared use by the VLM fallback.
_D7_FITZ_FALLBACK_ENABLED = os.getenv("D7_FITZ_FALLBACK_ENABLED", "true").strip().lower() in (
    "1",
    "true",
    "yes",
)


# ---------------------------------------------------------------------------
# libreoffice_to_pdf (lines 3769-3818)
# ---------------------------------------------------------------------------


def libreoffice_to_pdf(input_path: str) -> str:
    """Convert a DOCX/PPTX file to PDF via LibreOffice headless.

    Returns the path to the generated PDF in a temporary directory.
    The caller is responsible for cleaning up the parent directory:
        shutil.rmtree(os.path.dirname(pdf_path), ignore_errors=True)
    """
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        raise RuntimeError(
            "LibreOffice not found. Install libreoffice-headless and ensure it is on PATH."
        )
    outdir = tempfile.mkdtemp(prefix="lo_pdf_")
    # Each conversion gets its own profile dir so parallel invocations don't conflict.
    profile_dir = os.path.join(outdir, "lo_profile")
    os.makedirs(profile_dir, exist_ok=True)
    try:
        result = subprocess.run(
            [
                lo,
                f"-env:UserInstallation=file://{profile_dir}",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                outdir,
                input_path,
            ],
            capture_output=True,
            text=True,
            timeout=180,
        )
        stem = os.path.splitext(os.path.basename(input_path))[0]
        pdf_path = os.path.join(outdir, f"{stem}.pdf")
        # Check for the PDF first; a non-zero exit may be a recoverable warning
        if not os.path.isfile(pdf_path):
            pdfs = [f for f in os.listdir(outdir) if f.endswith(".pdf")]
            if pdfs:
                pdf_path = os.path.join(outdir, pdfs[0])
            elif result.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice conversion failed (exit {result.returncode}): "
                    f"{result.stderr.strip()}"
                )
            else:
                raise RuntimeError("LibreOffice did not produce a PDF file.")
        return pdf_path
    except Exception:
        shutil.rmtree(outdir, ignore_errors=True)
        raise


# ---------------------------------------------------------------------------
# html_to_markdown_with_images (lines 3821-3906)
# ---------------------------------------------------------------------------


async def html_to_markdown_with_images(path: str, model: str) -> str:
    """Convert an HTML file to markdown, replacing <img> tags with vision-API descriptions.

    Images are described concurrently via the OpenAI vision API and inserted as
    [Image: <description>] markers at the position of the original <img> tag.
    """
    import html2text

    with open(path, encoding="utf-8", errors="replace") as f:
        html_content = f.read()

    img_pattern = re.compile(r"<img[^>]+src=[\"']([^\"']+)[\"'][^>]*/?>", re.IGNORECASE)
    srcs = img_pattern.findall(html_content)

    async def _describe(src: str) -> str:
        import openai

        from ..client import get_openai_client
        from ..metrics import IMAGE_DESCRIBE_FAILURES

        async def _call() -> str:
            client = get_openai_client()
            response = await client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": src}},
                            {
                                "type": "text",
                                "text": (
                                    "Describe this image concisely in 1-2 "
                                    "sentences for document context."
                                ),
                            },
                        ],
                    }
                ],
                max_tokens=150,
            )
            return response.choices[0].message.content.strip()

        try:
            return await _call()
        except (openai.RateLimitError, openai.APIConnectionError):
            # Transient failure — retry once after a short backoff.
            await asyncio.sleep(2)
            try:
                return await _call()
            except openai.APIError as retry_exc:
                logger.error(
                    "Image description failed after retry (%s): %s",
                    type(retry_exc).__name__,
                    str(retry_exc)[:200],
                )
                IMAGE_DESCRIBE_FAILURES.labels(error_type=type(retry_exc).__name__).inc()
                return "image"
            except Exception:
                # Non-OpenAI error on retry (e.g. code bug) — do not swallow it.
                raise
        except openai.APIError as exc:
            logger.error(
                "Image description failed (%s): %s",
                type(exc).__name__,
                str(exc)[:200],
            )
            IMAGE_DESCRIBE_FAILURES.labels(error_type=type(exc).__name__).inc()
            return "image"

    descriptions = await asyncio.gather(*(_describe(src) for src in srcs))

    counter = iter(range(len(descriptions)))

    def _replace(match: re.Match) -> str:
        i = next(counter, None)
        desc = descriptions[i] if i is not None else "image"
        return f"[Image: {desc}]"

    modified_html = img_pattern.sub(_replace, html_content)

    h = html2text.HTML2Text()
    h.ignore_images = True
    h.ignore_links = False
    h.body_width = 0
    return normalize_dashes(h.handle(modified_html))


# ---------------------------------------------------------------------------
# xlsx_to_markdown (lines 3909-3943)
# ---------------------------------------------------------------------------


def xlsx_to_markdown(path: str) -> str:
    """Convert an .xlsx workbook to markdown tables (Fix 4; openpyxl is MIT, HR4).

    Each sheet becomes a `## <sheet>` title plus one markdown table (header row +
    separator + data rows), so the existing flat-table path (route_and_extract_flat ->
    _flat_parse_table) captures every cell into row_records and Fix-2 column-stitching
    is reused. Spreadsheets carry no heading hierarchy, so the document routes flat by
    design. read_only + data_only keeps memory bounded on large books."""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    out: list[str] = []
    try:
        for ws in wb.worksheets:
            rows = [
                ["" if c is None else str(c) for c in row] for row in ws.iter_rows(values_only=True)
            ]
            rows = [r for r in rows if any(cell.strip() for cell in r)]
            if not rows:
                continue
            width = max(len(r) for r in rows)
            out.append(f"## {ws.title}")
            header = rows[0] + [""] * (width - len(rows[0]))
            out.append("| " + " | ".join(c.replace("|", r"\|") for c in header) + " |")
            out.append("| " + " | ".join(["---"] * width) + " |")
            for r in rows[1:]:
                padded = r + [""] * (width - len(r))
                out.append("| " + " | ".join(c.replace("|", r"\|") for c in padded) + " |")
            out.append("")
    finally:
        wb.close()
    md = "\n".join(out).strip()
    if not md:
        raise RuntimeError(f"xlsx_to_markdown produced empty output for {path}")
    return md


# ---------------------------------------------------------------------------
# image_to_markdown (lines 3946-3963)
# ---------------------------------------------------------------------------


def image_to_markdown(path: str, ocr_lang_override: list[str] | None = None) -> str:
    """OCR a scanned image (.png/.jpg/.jpeg/.tiff) to markdown (Fix 4; OCR-only, HR3).

    An image has no text layer, so full-page OCR is always on. Routes the image through
    Docling's PDF/image pipeline with the same CPU-only Tesseract options as the PDF
    path (force_full_page_ocr + Fix-5 detected language). No VLM, no LLM egress -- local
    Tesseract only (HR3). VLM stays disabled by design (RFC-004)."""
    # Docling routes InputFormat.IMAGE through the same StandardPdfPipeline as PDF, so
    # this reuses the process-lifetime _docling_converter cache (see its docstring) --
    # a fresh DocumentConverter per call leaks ~237 MB RSS (torch/models never returned).
    converter = _docling_converter(
        force_full_page_ocr=True, ocr_lang_override=ocr_lang_override, for_image=True
    )
    result = converter.convert(path)
    md = _repair_docling_tables(result.document.export_to_markdown(), doc_name=path)
    if not md or not md.strip():
        raise RuntimeError(f"image_to_markdown produced empty output for {path}")
    return normalize_dashes(md)


# ---------------------------------------------------------------------------
# rasterize_pdf_pages (lines 3966-3989)
# ---------------------------------------------------------------------------


def rasterize_pdf_pages(pdf_path: str, dpi: int = 200) -> list[str]:
    """Rasterize each PDF page to a base64 data-URI PNG via pypdfium2 (HR4-safe)."""
    import base64
    import io

    import pypdfium2 as pdfium
    from PIL import Image

    pdoc = pdfium.PdfDocument(pdf_path)
    try:
        result: list[str] = []
        scale = dpi / 72
        for page_index in range(len(pdoc)):
            page = pdoc[page_index]
            bitmap = page.render(scale=scale)  # type: ignore[arg-type]
            pil_image: Image.Image = bitmap.to_pil()
            buf = io.BytesIO()
            pil_image.save(buf, format="PNG", optimize=True)
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            result.append(f"data:image/png;base64,{b64}")
            page.close()
        return result
    finally:
        pdoc.close()


# ---------------------------------------------------------------------------
# rasterize_pdf_pages_fitz (lines 4012-4045)
# ---------------------------------------------------------------------------


def rasterize_pdf_pages_fitz(pdf_path: str, dpi: int = 200) -> list[str]:
    """Rasterize each PDF page to a base64 data-URI PNG via fitz (D4, RFC-024).

    Fallback backend for ``rasterize_pdf_pages`` when pypdfium2 crashes on
    CMap-corrupt PDFs. Reuses the ``fitz.Page.get_pixmap()`` pattern already
    proven for image cropping in ``_recover_picture_text``."""
    from ..config import ALLOW_AGPL_FALLBACK

    if not ALLOW_AGPL_FALLBACK:
        # RFC-034 D4 step 4: this is a literal fallback backend, so a blocked
        # invocation is a prevented AGPL fallback — meter it for observability.
        from ..metrics import AGPL_FALLBACK_TOTAL

        AGPL_FALLBACK_TOTAL.labels(reason="blocked").inc()
        raise RuntimeError(
            f"cannot rasterize {pdf_path} via the fitz fallback backend: fitz "
            "(PyMuPDF, AGPL-3.0) is required and ALLOW_AGPL_FALLBACK=false"
        )
    import base64

    import fitz  # PyMuPDF, AGPL-3.0

    pdf = fitz.open(pdf_path)
    try:
        result: list[str] = []
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        for page in pdf:
            pix = page.get_pixmap(matrix=matrix)
            b64 = base64.b64encode(pix.tobytes("png")).decode("ascii")
            result.append(f"data:image/png;base64,{b64}")
        return result
    finally:
        pdf.close()


# ---------------------------------------------------------------------------
# tesseract_ocr_pdf_pages (lines 4048-4083)
# ---------------------------------------------------------------------------


async def tesseract_ocr_pdf_pages(pdf_path: str, langs: list[str]) -> str:
    """Rasterize each PDF page and OCR it via local Tesseract (RFC-023 D7).

    Last-resort recovery when the VLM fallback itself crashes on a garbled
    PDF -- no LLM egress, local ``tesseract`` binary only (HR3).

    D4 (RFC-024): tries pypdfium2 (``rasterize_pdf_pages``) first; on Exception
    (e.g. CMap-corrupt PDFs that crash pypdfium2), falls back to fitz
    (``rasterize_pdf_pages_fitz``) unless disabled via
    ``D7_FITZ_FALLBACK_ENABLED=false``."""
    import base64

    from .pictures import _tesseract_ocr_image

    try:
        page_images = await asyncio.to_thread(rasterize_pdf_pages, pdf_path)
    except Exception as exc:  # D4: fall back to fitz rasterization
        if not _D7_FITZ_FALLBACK_ENABLED:
            raise
        logger.warning(
            "rasterize_pdf_pages (pypdfium2) failed for %s (%s); falling back to fitz",
            pdf_path,
            exc,
        )
        page_images = await asyncio.to_thread(rasterize_pdf_pages_fitz, pdf_path)
    pages_text = []
    for data_uri in page_images:
        png_bytes = base64.b64decode(data_uri.split(",", 1)[1])
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as png_tmp:
            png_tmp.write(png_bytes)
            png_path = png_tmp.name
        try:
            text = await asyncio.to_thread(_tesseract_ocr_image, png_path, langs)
        finally:
            os.unlink(png_path)
        if text:
            pages_text.append(text)
    return "\n\n".join(pages_text)


# ---------------------------------------------------------------------------
# vlm_extract_markdown (lines 4086-4166)
# ---------------------------------------------------------------------------


async def vlm_extract_markdown(pdf_path: str, model: str | None = None) -> str:
    """Extract markdown from a PDF via a vision LLM — last-resort garble fallback."""
    import openai

    from ..client import get_openai_client
    from ..config import settings

    resolved_model = model or settings.vlm_model
    if resolved_model.startswith("azure/"):
        resolved_model = resolved_model[len("azure/") :]

    page_images = await asyncio.to_thread(rasterize_pdf_pages, pdf_path)
    if not page_images:
        raise RuntimeError(f"vlm_extract_markdown: no pages rasterized from {pdf_path}")

    client = get_openai_client()

    _VLM_PAGE_PROMPT = (
        "You are a document OCR assistant. Extract ALL visible text content from "
        "this scanned document page and return it as clean Markdown.\n\n"
        "Rules:\n"
        "- Preserve the document's heading hierarchy using Markdown heading levels "
        "(#, ##, ###, etc.).\n"
        "- Preserve tables as Markdown tables.\n"
        "- Preserve numbered and bulleted lists.\n"
        "- Ignore watermarks, background patterns, and page numbers.\n"
        "- If the page contains Arabic or right-to-left text, preserve the original "
        "script — do NOT transliterate.\n"
        "- Do NOT describe images; extract only text.\n"
        "- If the page is blank or contains no readable text, return exactly: "
        "<!-- blank page -->\n"
        "- Return ONLY the extracted Markdown, no commentary or wrapper."
    )

    async def _extract_page(page_idx: int, image_uri: str) -> tuple[int, str]:
        async def _call() -> str:
            response = await client.chat.completions.create(
                model=resolved_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": image_uri}},
                            {"type": "text", "text": _VLM_PAGE_PROMPT},
                        ],
                    }
                ],
                max_tokens=4096,
                temperature=0.0,
            )
            return response.choices[0].message.content.strip()

        try:
            return (page_idx, await _call())
        except (openai.RateLimitError, openai.APIConnectionError):
            await asyncio.sleep(2)
            try:
                return (page_idx, await _call())
            except Exception as retry_exc:
                logger.error("VLM page %d failed after retry: %s", page_idx + 1, retry_exc)
                return (page_idx, "")
        except Exception as exc:
            logger.error("VLM page %d extraction failed: %s", page_idx + 1, exc)
            return (page_idx, "")

    sem = asyncio.Semaphore(4)

    async def _bounded(idx: int, uri: str) -> tuple[int, str]:
        async with sem:
            return await _extract_page(idx, uri)

    results = await asyncio.gather(*[_bounded(i, u) for i, u in enumerate(page_images)])
    results_sorted = sorted(results, key=lambda r: r[0])
    page_markdowns = [md for _, md in results_sorted if md and md.strip() != "<!-- blank page -->"]

    if not page_markdowns:
        raise RuntimeError(
            f"vlm_extract_markdown: VLM returned no content for any page of {pdf_path}"
        )

    return "\n\n---\n\n".join(page_markdowns)


# ---------------------------------------------------------------------------
# docx_to_markdown (lines 4169-4190)
# ---------------------------------------------------------------------------


def docx_to_markdown(path: str) -> str:
    """Convert a DOCX file to a markdown string preserving heading hierarchy."""
    from docx import Document

    doc = Document(path)
    lines = []
    heading_map = {
        "Heading 1": "#",
        "Heading 2": "##",
        "Heading 3": "###",
        "Heading 4": "####",
        "Heading 5": "#####",
        "Heading 6": "######",
    }
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        prefix = next((v for k, v in heading_map.items() if para.style.name.startswith(k)), None)
        lines.append(f"{prefix} {text}" if prefix else text)
    return normalize_dashes("\n".join(lines))


# ---------------------------------------------------------------------------
# pptx_to_markdown (lines 4193-4213)
# ---------------------------------------------------------------------------


def pptx_to_markdown(path: str) -> str:
    """Convert a PPTX file to markdown, one H1 section per slide."""
    from pptx import Presentation

    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        title = (
            title_shape.text.strip() if title_shape and title_shape.text.strip() else f"Slide {i}"
        )
        lines.append(f"# {title}")
        for shape in slide.shapes:
            if shape == title_shape or not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
        lines.append("")
    return normalize_dashes("\n".join(lines))
