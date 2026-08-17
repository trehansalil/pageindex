"""Document format conversion helpers and tree search utilities."""

import asyncio
import contextlib
import dataclasses
import functools
import logging
import math
import multiprocessing
import os
import queue as queue_mod
import re
import shutil
import subprocess
import tempfile
import time
import unicodedata
from collections.abc import Callable
from concurrent.futures import ThreadPoolExecutor
from concurrent.futures import TimeoutError as FuturesTimeoutError
from typing import TYPE_CHECKING, TypedDict, cast

if TYPE_CHECKING:
    from docling.document_converter import DocumentConverter

from .picture_plane import OcrMode, PictureRegion, SkipReason, decide_ocr_mode
from .script import _AR_COMMON_WORDS as _AR_COMMON_WORDS
from .script import AR_CHAR_RE as _AR_LETTER_RE
from .script import AR_CHAR_RE as _AR_SCRIPT_RE
from .script import BlobKind, RtlDecision, _word_has_reversed_morphology, apply_rtl, decide_rtl, normalize_dashes
from .script import arabic_readability_score as _arabic_readability_score
from .script import is_arabic_char as _is_arabic_char

logger = logging.getLogger(__name__)


_HEADING_RE = re.compile(r"^(#{1,6})(?=\s)", re.MULTILINE)


def _relevel_headings(md: str) -> str:
    """Promote markdown headings so the shallowest present level becomes H1 (#)."""
    levels = [len(m.group(1)) for m in _HEADING_RE.finditer(md)]
    if not levels:
        return md
    shift = min(levels) - 1
    if shift <= 0:
        return md
    return _HEADING_RE.sub(lambda m: "#" * (len(m.group(1)) - shift), md)


# --- Heading-depth recovery from German-insurance numbering schemes -----------
# Docling's export_to_markdown() renders every section header at a single
# '#'-level, so even when the docling-hierarchical-pdf add-on selects the right
# headings the resulting tree is flat (depth 1) and fails the depth>=2 quality
# gate (HR5 / validate_tree). We re-derive each heading's depth from its
# numbering prefix. Two schemes appear in the validated German insurance corpus
# (2026-05-31):
#   dot notation   (e.g. AKB):  "A.1" -> 2, "A.1.1" -> 3   (bare "A" stays H1)
#   hyphen clauses (e.g. PHV):  "Abschnitt A1" -> 2, "A1-6" -> 3, "A1-6.1" -> 4
# numbering_depth() returns None when no scheme is recognised so generic /
# non-numbered documents keep their existing heading levels untouched.
_HLINE_RE = re.compile(r"^#{1,6}[ \t]+(.*\S)[ \t]*$", re.MULTILINE)
_NUM_SECTION_WORD_RE = re.compile(r"^Abschnitt\s+[A-Z]?\d", re.IGNORECASE)
_NUM_PART_RE = re.compile(r"^(?:Teil|Anhang|Kapitel|Abschnitt)\b", re.IGNORECASE)
_NUM_HYPHEN_RE = re.compile(r"^[A-Z]\d+-\d+(\.\d+)?(?=[ \t]|$)")
_NUM_DOT_RE = re.compile(r"^[A-Z](\.\d+)+(?=[ \t.:]|$)")
_NUM_PARA_RE = re.compile(r"^(?:§\s*)?\d+(\.\d+)+(?=[ \t.:]|$)")

# Arabic legal-gazette numbering (Fix 1 / RFC fizzy-forging-pearl). Arabic laws label
# structure with words, not Latin numbering: الباب/الفصل/القسم/الجزء (chapter/part/
# section) sit at the top; المادة ("Article", optionally bare مادة) is the per-article
# unit one level under. Matched on the Arabic stem (ال- prefix optional) so the depth
# signal is script-native — German/English titles never match these, so this tier is
# additive with zero regression risk on the existing corpus.
_AR_DIGITS = str.maketrans("٠١٢٣٤٥٦٧٨٩", "0123456789")
# RFC-033 D8: second alternative in each pattern matches the mirror-reversed
# (character-order-reversed) form Tesseract produces on some scanned inputs,
# e.g. 'المادة' -> 'ةداملا' (stem reversed, 'ال' prefix reversed to a trailing
# 'لا' suffix). Reversed stems: باب->باب (palindrome), فصل->لصف, قسم->مسق,
# جزء->ءزج, مادة->ةدام, مرسوم->موسرم. RFC-036 D5: same treatment for
# قرار->رارق, قانون->نوناق (gazette-style decree/law markers, part-level).
_AR_PART_RE = re.compile(
    r"^(?:ال)?(?:باب|فصل|قسم|جزء|قرار|مرسوم|قانون)\b"
    r"|^(?:باب|لصف|مسق|ءزج|رارق|موسرم|نوناق)(?:لا)?\b"
)
_AR_ARTICLE_RE = re.compile(r"^(?:ال)?مادة\b|^ةدام(?:لا)?\b")
# RFC-028 D1: char limit for wholesale heading promotion, raised from 60 to
# accommodate Arabic legal headings/titles (66-76+ chars observed).
_AR_HEADING_CHAR_LIMIT = 100
# RFC-028 D1: captures the structural marker word plus an immediately
# following PARENTHESIZED numeral (e.g. "مادة (3)") so a fused marker+title
# line exceeding the char limit can be split into a standalone heading (this
# capture) and remaining prose (everything after it). Deliberately requires
# the parenthetical — a bare trailing number ("مادة 2 من هذا القانون...") is
# the shape of a mid-paragraph/citation reference continuing into running
# prose ("Article 2 OF this law..."), not a title, so it must NOT be split
# into a heading (see TestInjectArabicStructuralHeadingsBlockStart's
# wrapped-citation case in test_rfc027_d4.py).
_AR_MARKER_CAPTURE_RE = re.compile(
    r"^(?:ال)?(?:باب|فصل|قسم|جزء|مادة|قرار|مرسوم|قانون)\s*\(\s*\d+\s*\)"
)

# Zone-3: _detect_arabic_reversal DELETED — replaced by decide_rtl (morphology+
# readability scorer) which is strictly stronger than the vocab-list method.
# _AR_KNOWN_WORDS, _AR_KNOWN_WORDS_REVERSED, _AR_REVERSAL_SAMPLE_THRESHOLD removed.


def _inject_arabic_structural_headings(md: str) -> str:
    """Promote al-bab/al-fasl/al-maddah lines to '#' headings (RFC-027 D4).

    Docling never classifies Arabic structural markers as ``SectionHeaderItem``s
    (English equivalents like "Chapter"/"Article" ARE detected), so these lines
    stay plain prose and the heading-depth recovery chain — which can only
    re-level EXISTING headings via ``_AR_PART_RE``/``_AR_ARTICLE_RE`` in
    ``numbering_depth`` — has nothing to work with, leaving the tree flat.

    Promotion is gated ONLY on the marker regex matching the line's **start**
    (RFC-028 D1) — a preceding-blank-line requirement is redundant/harmful for
    continuous OCR output, where scanned Arabic legal text flows without
    blank-line separators between consecutive مادة articles, and previously
    meant only the FIRST marker in a run was ever promoted. The anchor
    protects against mid-paragraph references like "...المشار إليها في
    المادة 5 من هذا القانون..." since those never start their line with the
    marker.

    A matching line up to 100 chars (RFC-028 D1: raised from 60 to
    accommodate Arabic legal headings like "المادة (3) نطاق التطبيق", which
    run 66-76+ chars) is promoted wholesale. A matching line that exceeds 100
    chars is split: the marker (plus any immediately-following numeral/
    parenthetical, e.g. "مادة (3)") becomes a standalone heading and the
    remaining prose is kept as a following line rather than dropped. Depth is
    left to the existing ``_relevel_by_containment``/``_relevel_by_numbering``
    chain.

    RFC-033 D8: when ``decide_rtl`` judges ``md`` to be Tesseract
    mirror-reversed, each line is character-reversed before pattern
    matching (the regexes are forward-oriented) but the ORIGINAL line text —
    whatever Tesseract actually produced — is what gets promoted to a
    heading; only the matching step operates on the flipped text."""
    reversed_ocr = decide_rtl(md).reversed
    lines = md.split("\n")
    out = []
    for line in lines:
        t = line.strip()
        if t and not _HEADING_RE.match(t):
            match_t = t[::-1] if reversed_ocr else t
            is_part = bool(_AR_PART_RE.match(match_t))
            is_article = is_part is False and bool(_AR_ARTICLE_RE.match(match_t))
            if is_part or is_article:
                level = "#" if is_part else "##"
                if len(t) <= _AR_HEADING_CHAR_LIMIT:
                    out.append(f"{level} {t}")
                    continue
                marker_match = _AR_MARKER_CAPTURE_RE.match(match_t)
                if marker_match:
                    if reversed_ocr:
                        marker_len = marker_match.end()
                        marker = t[len(t) - marker_len :]
                        remainder = t[: len(t) - marker_len].strip()
                    else:
                        marker = marker_match.group(0).strip()
                        remainder = t[marker_match.end() :].strip()
                    out.append(f"{level} {marker}")
                    if remainder:
                        out.append(remainder)
                    continue
        out.append(line)
    return "\n".join(out)


_DE_ZIFFER_RE = re.compile(r"^(?:Ziffer|Ziff\.)\s+\d+")

# Line-start anchoring alone does not stop a *clause body* that opens with its
# own number ("Ziffer 3 gilt entsprechend für ...", "Article (5) shall apply
# where ...") from being promoted, which would swallow a whole paragraph into a
# heading title. Real clause titles are short; the limit is deliberately well
# above the longest realistic one (German AHB titles such as "Ziffer 7
# Versicherungsfall, Obliegenheiten des Versicherungsnehmers ..." run past 100
# chars, so the Arabic function's 100-char limit would clip them) and only
# blocks obvious running prose. Over-limit lines are left untouched — nothing
# is split or dropped.
_CLAUSE_HEADING_CHAR_LIMIT = 200


def _inject_german_clause_headings(md: str) -> str:
    """Promote 'Ziffer N'/'Ziff. N' lines to '##' headings (RFC-033 D5).

    German AHB/insurance documents use Ziffer/Ziff. clause numbering that
    Docling does not detect as headings. The line-start anchor prevents
    mid-sentence references like "see Ziffer 1 above" from being promoted, and
    ``_CLAUSE_HEADING_CHAR_LIMIT`` prevents a paragraph that merely opens with
    a clause reference from becoming a heading."""
    lines = md.split("\n")
    out = []
    for line in lines:
        t = line.strip()
        if (
            t
            and len(t) <= _CLAUSE_HEADING_CHAR_LIMIT
            and not _HEADING_RE.match(t)
            and _DE_ZIFFER_RE.match(t)
        ):
            out.append(f"## {t}")
            continue
        out.append(line)
    return "\n".join(out)


_EN_ARTICLE_PROSE_RE = re.compile(r"^Article\s*\(\s*\d+\s*\)", re.IGNORECASE)


def _inject_english_article_headings(md: str) -> str:
    """Promote 'Article (N)' lines to '##' headings (RFC-033 D5).

    UAE/English legal documents use parenthesised 'Article (N)' numbering
    that Docling sometimes misses entirely as a heading, leaving it as plain
    prose. The line-start anchor prevents mid-sentence references like "see
    Article (1) above" from being promoted, and ``_CLAUSE_HEADING_CHAR_LIMIT``
    prevents a paragraph that merely opens with an article reference from
    becoming a heading."""
    lines = md.split("\n")
    out = []
    for line in lines:
        t = line.strip()
        if (
            t
            and len(t) <= _CLAUSE_HEADING_CHAR_LIMIT
            and not _HEADING_RE.match(t)
            and _EN_ARTICLE_PROSE_RE.match(t)
        ):
            out.append(f"## {t}")
            continue
        out.append(line)
    return "\n".join(out)


def numbering_depth(title: str) -> int | None:
    """Infer a 1-based heading depth from a German-insurance numbering prefix.

    Returns None when no recognised numbering scheme is present, so the caller
    leaves such headings at their existing level."""
    t = title.strip()
    # Arabic structural words (Fix 1): chapter/part -> top, Article -> one level under.
    if _AR_PART_RE.match(t):
        return 1
    if _AR_ARTICLE_RE.match(t):
        return 2
    # "Abschnitt A1" is a section nested one level under its "Teil".
    if _NUM_SECTION_WORD_RE.match(t):
        return 2
    # Part / appendix words sit at the top.
    if _NUM_PART_RE.match(t):
        return 1
    # Hyphen clauses: "A1-6" -> 3, "A1-6.1" -> 4.
    if _NUM_HYPHEN_RE.match(t):
        return 4 if "." in t.split()[0] else 3
    # Dot notation: "A.1" -> 2, "A.1.1" -> 3.
    m = _NUM_DOT_RE.match(t)
    if m:
        return 1 + m.group(0).count(".")
    # Plain paragraph / numeric sub-sections: "3.1" -> 2, "§3.1.2" -> 3.
    m = _NUM_PARA_RE.match(t)
    if m:
        return 1 + m.group(0).count(".")
    return None


def _relevel_by_numbering(md: str) -> str:
    """Override each markdown heading's '#'-level from its numbering prefix.

    Headings whose title has no recognised numbering prefix are left unchanged,
    so this is safe to run after ``_relevel_headings`` on any corpus."""

    def repl(m: "re.Match[str]") -> str:
        title = m.group(1)
        depth = numbering_depth(title)
        if depth is None:
            return m.group(0)
        # Clamp to the markdown heading range [1, 6]; deeply nested numbered
        # sections (e.g. "A.1.1.1.1.1") would otherwise emit 7+ '#'s, which is
        # not a valid heading. Mirrors the clamp in _relevel_by_containment.
        return "#" * max(1, min(6, depth)) + " " + title

    return _HLINE_RE.sub(repl, md)


# --- Heading-depth recovery by numbering-prefix CONTAINMENT (no per-scheme table) -
# Every STRUCTURAL depth signal Docling exposes on this corpus is flat:
# SectionHeaderItem.level, body-tree traversal depth, and the PDF outline are all
# level 1 (verified 2026-05-31). So depth must be inferred from the heading TEXT.
# We segment each heading's leading numbering LABEL into atomic components
# ("A.1.1"->[A,1,1]; "Abschnitt A1"->[A,1]; bare prose title -> []) and set
# depth = 1 + length of the longest OTHER present label that is a proper prefix.
# This nests an unseen numbering style without a hardcoded regex table, so it is
# the PRIMARY depth source; numbering_depth() above is kept only as a last-resort
# fallback for the degenerate case where containment stays flat.

# Structural words that introduce a numbering label. We collapse the spaced-out
# Docling rendering ("T e i l   A") before matching, so the regex sees "TeilA".
_WORD_RE = re.compile(r"^(teil|anhang|abschnitt|kapitel)\b", re.IGNORECASE)
# Arabic structural words consumed the same way as the Latin _WORD_RE (Fix 1): the
# label that NESTS is the number/letter that follows المادة/الباب/الفصل/القسم/الجزء/مرسوم.
# RFC-033 D8: reversed-form alternative, same construction as _AR_PART_RE/
# _AR_ARTICLE_RE above.
_AR_WORD_RE = re.compile(
    r"^(?:ال)?(?:مادة|باب|فصل|قسم|جزء|مرسوم)\b|^(?:ةدام|باب|لصف|مسق|ءزج|موسرم)(?:لا)?\b"
)
# English "Article N" / "Art. N" and section-symbol "§ N" headings (RFC-015 D3
# Part B / task 1.4): without this, _segment_label() rejects "Article 5" as
# prose (its first alnum component is the multi-letter word "Article", which
# fails the single-letter-or-digit label-first-component gate below) and
# returns [] — an unlabelled/"bare title" heading whose ORIGINAL doc-derived
# level is left untouched by _relevel_by_containment. On corpora where that
# original level is deeply nested under an unrelated sub-bullet, Articles 3-6
# inherit that nesting verbatim ("staircase" mis-nesting) instead of getting an
# explicit depth. Recognising the "Art(icle|.) N" / "§ N" prefix as a
# structural word (consumed the same way as Teil/Abschnitt/Kapitel above)
# lets the bare number become the label, so containment assigns depth 1.
# The number may also be parenthesised — "Article (47)", "§ (12)" — as used by
# UAE/English legal corpora (RFC-033 D4); the optional '(' is matched here and
# stripped downstream by _segment_label's head.strip("()").
_ARTICLE_RE = re.compile(r"^(?:Art(?:icle|\.)\s+\(?\s*\d+|§\s*\(?\s*\d+)", re.IGNORECASE)
_ARTICLE_WORD_RE = re.compile(r"^(?:Art(?:icle|\.)|§)\s*", re.IGNORECASE)


def _collapse_spaced(text: str) -> str:
    """Collapse Docling's letter-spaced headings: 'T e i l   A' -> 'Teil A'.

    Docling renders these with SINGLE spaces between the letters of a word and a
    WIDER gap (2+ spaces) between words. We split on runs of 2+ spaces to recover
    word boundaries, then glue single-spaced letters inside each chunk. A '-'
    surrounded by spaces is kept as a separator word. Ordinary headings (whose
    tokens are multi-letter) are returned unchanged.
    """
    raw_toks = text.split()
    if not (len(raw_toks) >= 4 and sum(1 for t in raw_toks if len(t) == 1) >= len(raw_toks) * 0.6):
        return text
    # split on 2+ spaces -> word-level chunks; within a chunk, single chars glue.
    chunks = re.split(r"\s{2,}", text.strip())
    out = []
    for chunk in chunks:
        ctoks = chunk.split()
        buf = []
        for t in ctoks:
            if t == "-":
                # A '-' inside a chunk is a label separator (e.g. clause code
                # "A1-6.1" rendered spaced as "A 1 - 6 . 1"): glue it to the
                # surrounding letters with NO spaces so the clause code stays
                # intact. Surrounding it with spaces would break hyphenated
                # clause-code detection ("A1-6.1" -> "A1- 6.1").
                buf.append("-")
            else:
                buf.append(t)
        if buf:
            out.append("".join(buf))
    return " ".join(out)


def _split_alnum(tok: str) -> list[str]:
    """Split an alnum label token at every letter<->digit boundary and (..) group.

    "A1"   -> ["A","1"]      "A(GB)" -> ["A","GB"]      "B4" -> ["B","4"]
    "A"    -> ["A"]          "A(GB)1"-> ["A","GB","1"]
    """
    # pull out a parenthesised group first
    parts: list[str] = []
    m = re.match(r"^([A-Za-z]+)(?:\(([A-Za-z0-9]+)\))?(\d+)?$", tok)
    if m:
        if m.group(1):
            parts.append(m.group(1))
        if m.group(2):
            parts.append(m.group(2))
        if m.group(3):
            parts.append(m.group(3))
        return parts
    # fallback: generic letter/digit run split
    return list(re.findall(r"[A-Za-z]+|\d+", tok))


def _segment_label(title: str) -> list[str]:
    """Segment a heading's leading numbering label into atomic components.

    Returns [] when the heading carries no recognisable label (a bare title).
    Word-prefix rule: a leading structural word (Teil/Anhang/Abschnitt/Kapitel)
    is consumed; the label that follows it is what nests. "Teil A"->[A];
    "Abschnitt A1"->[A,1]; "A1-6.1"->[A,1,6,1]; "A.1.1"->[A,1,1];
    "A(GB)-1"->[A,GB,1]; "Versicherte Personen"->[].
    """
    t = _collapse_spaced(title.strip())
    # consume an optional leading structural word (Latin or Arabic, Fix 1)
    wm = _WORD_RE.match(t)
    if wm:
        t = t[wm.end() :].lstrip(" -")
    elif _ARTICLE_RE.match(t):
        # "Article 5" / "Art. 5" / "§ 12" (RFC-015 D3 Part B): consume just the
        # word/symbol, leaving the bare number as the label so it gets an
        # explicit containment depth instead of falling through as prose.
        t = _ARTICLE_WORD_RE.sub("", t, count=1)
    else:
        awm = _AR_WORD_RE.match(t)
        if awm:
            t = t[awm.end() :].lstrip(" -:()")
    # Normalise Arabic-Indic digits so an Arabic article number ("المادة ٩") is
    # recognised by the same digit-aware label logic as Latin numbering. No-op on
    # ASCII, so German/English labels are byte-for-byte unaffected.
    t = t.translate(_AR_DIGITS)
    # the label is the leading run of [alnum . - ( )] up to the first space
    # that starts the descriptive title. Grab the first whitespace-delimited tok.
    head = t.split(maxsplit=1)[0] if t else ""
    if not head:
        return []
    # The label may itself contain '.', '-' and '()' separators.
    # Stop the label at a separator that is followed by a non-label char? Simpler:
    # the head token IS the label candidate. Validate it starts with a letter
    # and contains at least one alnum.
    # Strip trailing punctuation like ':' or '.' used as a terminator? Keep dots
    # that are internal (A.1.) — drop a single trailing '.'/':'.
    head = head.rstrip(":")
    # Strip wrapping parentheses on the whole token ("(9)" -> "9") so a parenthesised
    # Arabic article number survives the leading-char gate. Internal parens of a Latin
    # label like "A(GB)1" are untouched (no leading/trailing paren on the token).
    head = head.strip("()")
    # A label must begin with a letter (clause code "A1-6.1") OR a digit
    # (numeric section "3.1") to be a recognisable numbering label. Rejecting
    # digit-led heads here would drop numeric headings before the digit-aware
    # validation below, leaving that part of the hierarchy flat.
    if not re.match(r"^[A-Za-z0-9]", head):
        return []
    comps: list[str] = []
    # split on the structural separators '.', '-'
    for seg in re.split(r"[.\-]", head):
        seg = seg.strip()
        if not seg:
            continue
        sub = _split_alnum(seg)
        if not sub:
            return []  # contains a non-alnum chunk we don't understand -> no label
        comps.extend(sub)
    # Reject degenerate labels: a single letter that is actually a word start
    # is fine ("A"), but require the whole head to be alnum/sep only — if the
    # head had spaces stripped we already isolated one token, so this holds.
    # Guard: a pure single-letter label is valid (top section).
    if not comps:
        return []
    # Reject labels that are clearly prose (e.g. first word "Was", "Wer"): a real
    # label is short and its first component is a single letter OR is all digits.
    first = comps[0]
    if not (re.fullmatch(r"[A-Za-z]", first) or first.isdigit()):
        return []
    # And the whole token must be short-ish (a clause code, not a sentence word).
    if len(head) > 14:
        return []
    return comps


def _containment_depths(titles: list[str]) -> list[int | None]:
    """Depth of each heading via numbering-prefix containment (grammar inference).

    depth(i) = 1 + length of the longest OTHER label that is a proper prefix of
    label(i). Bare-title headings (label == []) return None so the caller leaves
    that heading's existing level untouched."""
    labels = [_segment_label(t) for t in titles]
    label_set = [tuple(lbl) for lbl in labels]
    present = {lbl for lbl in label_set if lbl}
    depths: list[int | None] = []
    for lab in label_set:
        if not lab:
            depths.append(None)
            continue
        # longest proper prefix that is itself a present label
        best = 0
        for k in range(len(lab) - 1, 0, -1):
            if lab[:k] in present:
                best = k
                break
        depths.append(best + 1)
    return depths


def _relevel_by_containment(md: str) -> str:
    """Override each markdown heading's '#'-level from numbering-prefix containment.

    Headings whose title carries no label (containment depth None) are left
    exactly as-is, so this is safe to run after ``_relevel_headings`` on any
    corpus. Non-heading text and spacing are preserved verbatim."""
    matches = list(_HLINE_RE.finditer(md))
    if not matches:
        return md
    depths = _containment_depths([m.group(1) for m in matches])
    out: list[str] = []
    pos = 0
    for m, depth in zip(matches, depths, strict=False):
        out.append(md[pos : m.start()])
        if depth is None:
            out.append(m.group(0))  # no label -> keep existing level
        else:
            out.append("#" * max(1, min(6, depth)) + " " + m.group(1))
        pos = m.end()
    out.append(md[pos:])
    return "".join(out)


def _max_heading_level(md: str) -> int:
    """Largest markdown heading '#'-run length present, or 0 if none."""
    levels = [len(m.group(1)) for m in _HEADING_RE.finditer(md)]
    return max(levels) if levels else 0


# --- Heading-depth recovery from the PDF OUTLINE (last resort for flat prose) --
# German IPID / FAQ insurance PDFs (Katzen-/Hunde-/Pferde-Kranken/-OP, Meuten,
# Halterhaftpflicht) have headings that are ALL bare prose with no numbering
# prefix, so _relevel_by_containment and _relevel_by_numbering both leave every
# heading at H1 (max_heading_level == 1) -> a FALSE depth<2 rejection of a
# legitimately-structured document. Their only author-declared structure is the
# PDF outline (PyMuPDF get_toc). We map each rendered heading to the outline
# section its PAGE falls in (text matching alone is insufficient: a TOC title such
# as "Informationsblatt zu Versicherungsprodukten" is often NOT in Docling's
# rendered heading set, so it must be located/injected by PAGE, never by text).
# A rendered heading whose text matches its section title becomes that section's
# anchor (H{toc_level}); every other heading in the section becomes a child
# (H{deepest_covering_level + 1}); a covered section title Docling never rendered
# is injected verbatim at its declared level. Cat A numbered docs never reach this
# step (the guard requires max_heading_level < 2 after the numbering chain); Cat D
# leaflets have no usable outline (<2 entries) and stay a legitimate depth<2
# rejection (HR5 — the quality gate is never weakened).
_OUTLINE_MIN_ANCHOR_ALNUM = (
    8  # min alnum chars for a substring title match (avoids short-title false positives)
)


def _outline_norm(s: str) -> str:
    """Normalise a title to lowercase alphanumerics for cross-source matching.

    Unifies dashes, collapses embedded newlines, then strips to [a-z0-9] — the
    same normalisation idiom the add-on's infer() / _patch_hierarchical_infer use,
    so a PyMuPDF outline title reconciles with a Docling-rendered heading despite
    whitespace, dash-variant and punctuation differences."""
    return re.sub(r"[^a-z0-9]", "", normalize_dashes((s or "").replace("\n", " ")).lower())


def _title_matches(norm_heading: str, norm_section: str) -> bool:
    """True when a rendered heading IS its outline section's title.

    Exact normalised equality, or a substring match when the shorter string is
    substantial (>= _OUTLINE_MIN_ANCHOR_ALNUM alnum chars) — this tolerates
    Docling rendering a longer heading than the TOC title (or vice versa) without
    the short-title false positives a bare endswith/startswith would admit."""
    if not norm_heading or not norm_section:
        return False
    if norm_heading == norm_section:
        return True
    shorter, longer = (
        (norm_heading, norm_section)
        if len(norm_heading) <= len(norm_section)
        else (norm_section, norm_heading)
    )
    return len(shorter) >= _OUTLINE_MIN_ANCHOR_ALNUM and shorter in longer


def _collect_heading_pages(doc) -> dict[str, list[int]]:
    """Map normalised-heading-text -> [page_no, ...] from a Docling document.

    Page numbers are 1-indexed (pypdfium2 backend: prov[0].page_no == page index
    + 1), matching PyMuPDF get_toc's page field. Pages are appended in document
    iteration order so repeated identical headings (e.g. the 3x "Besondere
    Bedingungen ..." chapters in Hundehalterhaftpflicht) can be disambiguated by a
    consumption pointer in _apply_outline_levels. Call this on the document state
    matching the markdown the caller will relevel (the RAW pre-add-on document on
    the over-prune raw_md fallback path; the post-add-on document otherwise) so the
    page map and the heading set stay in sync."""
    from collections import defaultdict

    from docling_core.types.doc.document import SectionHeaderItem

    pages: dict[str, list[int]] = defaultdict(list)
    for item, _ in doc.iterate_items(with_groups=False):
        if not isinstance(item, SectionHeaderItem) or not item.prov:
            continue
        key = _outline_norm(item.text or "")
        if key:
            pages[key].append(item.prov[0].page_no)
    return dict(pages)


def _read_pdf_outline(pdf_path: str) -> tuple[list[tuple[int, str, int]], int]:
    """Read the PDF bookmark outline as [(level, title, page_1indexed), ...] in
    document (outline-tree) order, plus the total page count.

    Uses pypdfium2 (BSD-3/Apache-2), NOT PyMuPDF (AGPL-3.0, HR4): this is the only
    first-party structural read on the default / VLM-bound PDF path, so keeping it
    off AGPL means zero first-party AGPL code touches a document the gate may later
    escalate (RFC-004 Q2). pypdfium2 reports level and page index 0-based; we add
    +1 to each to preserve the PyMuPDF get_toc convention the pure consumer
    (_apply_outline_levels) expects. The +1 on level is LOAD-BEARING: a 0-based
    level would collapse depth-1 and depth-2 sections through max(1, min(6, level)).
    A bookmark with no resolvable page destination maps to page 0 (a sentinel that
    never falls inside a 1-indexed section's [start, end) range).

    Returns ([], 0) when the outline has fewer than 2 entries — no usable
    structural signal, so the caller leaves the markdown flat and the quality gate
    rejects it legitimately (Cat D leaflets). Document order is preserved (NOT
    sorted by page): section extents are computed by outline NESTING (the next
    entry whose level <= the current level), which requires reading order."""
    import pypdfium2 as pdfium

    pdoc = pdfium.PdfDocument(pdf_path)
    try:
        toc: list[tuple[int, str, int]] = []
        for bm in pdoc.get_toc():
            dest = bm.get_dest()
            page_index = dest.get_index() if dest is not None else None
            page_1based = (page_index + 1) if page_index is not None else 0
            toc.append((bm.level + 1, bm.get_title() or "", page_1based))
        total_pages = len(pdoc)
    finally:
        pdoc.close()
    if len(toc) < 2:
        return [], 0
    return toc, total_pages


# Complexity grandfathered (outline relevel, depth<2 fix); see pyproject [tool.ruff].
def _apply_outline_levels(  # noqa: C901, PLR0915
    md: str,
    heading_pages: dict[str, list[int]],
    toc: list[tuple[int, str, int]],
    total_pages: int,
) -> str:
    """PURE last-resort relevel: assign each markdown heading an H-level from the
    PDF-outline section its page falls in, injecting any outline section title that
    Docling never rendered. No Docling/PyMuPDF deps -> directly unit-testable.

    Inputs:
      heading_pages : {normalised_title -> [page_no, ...]} in document order
      toc           : [(level, title, page_1indexed), ...] in outline order
      total_pages   : page count (bounds the last section's extent)

    Section extents follow outline NESTING: section i spans [page_i, next_start)
    where next_start is the page of the next entry whose level <= level_i (or
    total_pages + 1 for the last). A page may be covered by several nested sections
    (a parent and its child); a heading is assigned its DEEPEST covering section's
    level + 1 (a child), unless its text matches one of its covering section titles
    (an anchor -> that section's level). A covered section whose title no rendered
    heading matched has its title injected, verbatim, at its declared level before
    the first rendered heading in its range.

    Returns md unchanged when there is no usable outline / no headings, or when the
    rewrite still yields depth<2 (so the gate rejects it rather than receiving a
    worse tree — HR5)."""
    if not toc:
        return md
    matches = list(_HLINE_RE.finditer(md))
    if not matches:
        return md

    # 1. Sections as half-open page ranges [start, end) respecting outline nesting.
    n = len(toc)
    sections: list[dict] = []
    for i, (level, title, start) in enumerate(toc):
        end = total_pages + 1
        for j in range(i + 1, n):
            if toc[j][0] <= level:
                end = toc[j][2]
                break
        sections.append(
            {
                "level": max(1, min(6, level)),
                "norm": _outline_norm(title),
                "raw": re.sub(r"\s+", " ", normalize_dashes(title)).strip(),
                "start": start,
                "end": end,  # exclusive
                "matched": False,  # a rendered heading was this section's title
            }
        )

    # 2. Consumption pointer for repeated identical headings (document order).
    from collections import deque

    page_q: dict[str, deque] = {k: deque(v) for k, v in heading_pages.items()}

    def _pop_page(norm_title: str) -> int | None:
        q = page_q.get(norm_title)
        return q.popleft() if q else None

    # 3. Assign a level to every rendered heading; record the first rendered
    #    heading offset per covering section (the injection insertion point).
    new_levels: list[int | None] = []
    first_off: dict[int, int] = {}
    for m in matches:
        norm_h = _outline_norm(m.group(1))
        page_no = _pop_page(norm_h)
        if page_no is None:
            new_levels.append(None)  # no provenance -> leave at current level
            continue
        covering = [(idx, s) for idx, s in enumerate(sections) if s["start"] <= page_no < s["end"]]
        if not covering:
            new_levels.append(None)  # cover/frontmatter before the first section
            continue
        for idx, _s in covering:
            first_off.setdefault(idx, m.start())
        # anchor = deepest covering section whose title this heading matches
        anchor = None
        for _idx, s in sorted(covering, key=lambda c: -c[1]["level"]):
            if _title_matches(norm_h, s["norm"]):
                anchor = s
                break
        if anchor is not None:
            new_levels.append(anchor["level"])
            anchor["matched"] = True
        else:
            deepest = max(covering, key=lambda c: c[1]["level"])[1]
            new_levels.append(max(1, min(6, deepest["level"] + 1)))

    # 4. Inject section titles Docling never rendered (e.g. the IPID overview).
    injections: dict[int, list[tuple[int, str]]] = {}
    for idx, s in enumerate(sections):
        if s["matched"]:
            continue
        off = first_off.get(idx)
        if off is None:
            continue  # no rendered heading in range -> nothing to anchor to
        line = "#" * s["level"] + " " + s["raw"] + "\n"
        injections.setdefault(off, []).append((s["level"], line))

    # 5. Splice: rewrite levels + emit injected titles (shallowest level first).
    out: list[str] = []
    pos = 0
    for m, lvl in zip(matches, new_levels, strict=False):
        out.append(md[pos : m.start()])
        if m.start() in injections:
            for _lvl, line in sorted(injections[m.start()]):
                out.append(line)
        out.append(m.group(0) if lvl is None else "#" * lvl + " " + m.group(1))
        pos = m.end()
    out.append(md[pos:])
    result_md = "".join(out)

    if _max_heading_level(result_md) < 2:
        return md  # outline gave no usable depth -> stay a legitimate rejection
    return result_md


def _relevel_by_outline(md: str, heading_pages: dict[str, list[int]], pdf_path: str) -> str:
    """Thin I/O wrapper: read the PDF outline, then apply outline-derived levels.

    Last-resort depth recovery for flat-prose docs after the numbering chain fails.
    Never fatal (the caller also wraps it); returns md unchanged on any failure or
    when the outline has <2 entries (Cat D leaflet)."""
    toc, total_pages = _read_pdf_outline(pdf_path)
    if not toc:
        return md
    result_md = _apply_outline_levels(md, heading_pages, toc, total_pages)
    logger.info(
        "_relevel_by_outline: applied outline page-spine to %s (%d sections, max_level %d)",
        pdf_path,
        len(toc),
        _max_heading_level(result_md),
    )
    return result_md


def _splice_landscape_fallback(
    md: str, landscape_fallback_pages: list[dict], heading_pages: dict[str, list[int]]
) -> str:
    """RFC-036 D0d: insert landscape-fallback markdown at its original page
    position in the document's block sequence, instead of appending after the
    last block.

    Insertion offset = the start of the first heading (in ``md``, matched to a
    page via ``heading_pages`` the same way ``_apply_outline_levels`` does)
    whose page number exceeds the fallback page's ``page_no`` — i.e. the next
    heading following that page. Falls back to document end when no such
    heading exists (the fallback page is on/after the document's last heading),
    which reproduces the prior append-at-end behaviour for that case. A no-op
    when ``landscape_fallback_pages`` is empty — documents that never trigger
    the landscape path see no ordering change."""
    if not landscape_fallback_pages:
        return md
    from collections import deque

    matches = list(_HLINE_RE.finditer(md))
    page_q: dict[str, deque] = {k: deque(v) for k, v in heading_pages.items()}
    heading_offsets: list[tuple[int, int | None]] = []
    for m in matches:
        norm_h = _outline_norm(m.group(1))
        q = page_q.get(norm_h)
        heading_offsets.append((m.start(), q.popleft() if q else None))

    inserts: list[tuple[int, int, str]] = []
    for p in landscape_fallback_pages:
        block = p["markdown"].strip()
        if not block:
            continue
        # Fallback page_no is PyMuPDF 0-indexed (_tag_landscape_pages_for_fallback);
        # heading_pages values are Docling prov.page_no, 1-indexed — same
        # correction as _landscape_pages_below_threshold.
        fallback_page_1idx = p["page_no"] + 1
        insert_at = len(md)
        for offset, page_no in heading_offsets:
            if page_no is not None and page_no > fallback_page_1idx:
                insert_at = offset
                break
        inserts.append((insert_at, fallback_page_1idx, block))

    # Descending offset so earlier insertions don't shift later ones; the
    # secondary page_no key keeps same-offset blocks in ascending page order
    # in the final document (last-inserted lands first at a shared offset).
    for insert_at, _page, block in sorted(inserts, reverse=True):
        md = md[:insert_at] + "\n\n" + block + "\n\n" + md[insert_at:]
    return md


def _has_structural_depth(md: str) -> bool:
    """Structural proxy for validate_tree's ``node_count>=3 AND depth>=2`` — used
    only to SELECT the better markdown SOURCE; the real HR5 gate still runs in the
    client. md_to_tree makes one tree node per heading and nests by '#'-level, so a
    depth>=2 tree needs a heading at level>=2 and node_count>=3 needs >=3 headings.
    Conservative by design: a false 'pass' here is still caught by the real gate."""
    return _max_heading_level(md) >= 2 and len(_HEADING_RE.findall(md)) >= 3


def _recover_heading_depth(md: str, heading_pages: dict[str, list[int]], pdf_path: str) -> str:
    """Run the full heading-depth recovery chain on ONE markdown source.

    containment (PRIMARY numbering-prefix depth) -> numbering (per-scheme regex
    FALLBACK) -> PDF outline (LAST RESORT for numberless flat prose). Each step
    runs only if the prior left the tree degenerately flat (max_heading_level<2).
    Dashes are normalised first so hyphen clause codes (A1-6.1) parse. The outline
    step is wrapped so it is never fatal."""
    md = _relevel_by_containment(_relevel_headings(normalize_dashes(md)))
    if _max_heading_level(md) < 2:
        md = _relevel_by_numbering(md)
    if _max_heading_level(md) < 2:
        try:
            md = _relevel_by_outline(md, heading_pages, pdf_path)
        except Exception as exc:
            logger.warning(
                "_relevel_by_outline failed for %s (%s); leaving flat markdown",
                pdf_path,
                exc,
            )
    return md


def _candidate_from_document(
    md: str, heading_pages: dict[str, list[int]], pdf_path: str
) -> "Candidate":
    """Build and depth-recover a single pipeline candidate.

    Runs ``_build_candidate`` (injection + normalisation) then
    ``_recover_heading_depth`` (containment -> numbering -> outline) and
    returns an immutable ``Candidate`` bundling the result with its heading-
    page map.  Pure function — does NOT call ``_collect_heading_pages`` (the
    caller owns that) and does NOT touch ``extraction_stages``.
    """
    built = _build_candidate(md)
    recovered = _recover_heading_depth(built, heading_pages, pdf_path)
    return Candidate(md=recovered, heading_pages=heading_pages)


# RFC-015 D5d: a valid sub-clause component is a digit run with an OPTIONAL single
# lowercase-letter suffix ("10a"), OR a lone lowercase letter ("a"). The lone-letter
# alternative is load-bearing: _segment_label() splits a letter suffix into its OWN
# component ("A.1.1a" -> [A,1,1,a]; "§ 5a" -> [5,a]), so the blueprint's literal
# `\d+[a-z]?` alone would never match the standalone trailing "a" in its own worked
# example ('7','10','a'). Widening to `\d+[a-z]?|[a-z]` promotes letter-suffixed
# sub-clauses (doc acc20e08) while a BARE list marker "a" (label ["a"], no numeric
# anchor prefix) still cannot promote — the k-loop below requires a non-empty prefix.
_SUBCLAUSE_COMP_RE = re.compile(r"\d+[a-z]?|[a-z]")


def _is_numeric_extension(lab: tuple, anchors: set) -> bool:
    """True when ``lab`` is a kept anchor label plus a run of numeric / letter-suffix
    sub-components (RFC-015 D5d).

    There must be a non-empty kept-section label P (in ``anchors``) that is a PROPER
    prefix of ``lab``, and every component beyond P must be a pure digit run, a digit
    run with a single trailing lowercase letter, or a lone lowercase letter. The
    proper-prefix requirement (``k`` never reaches 0) means a bare list marker cannot
    promote itself."""
    return any(
        lab[:k] in anchors and all(_SUBCLAUSE_COMP_RE.fullmatch(c) for c in lab[k:])
        for k in range(len(lab) - 1, 0, -1)
    )


def _repromote_numbered_headings(doc) -> int:
    """Re-promote demoted body TextItems back to headings (no hardcoding).

    The docling-hierarchical-pdf add-on gives a clean heading SELECTION but
    over-prunes: it demotes deep numbered clauses (AKB "A.1.1", "A.1.1.1") to
    body TextItems alongside the font-size junk, capping the tree's depth. This
    walks the post-add-on doc and converts a TextItem back to a SectionHeaderItem
    IFF its numbering label is a proper NUMERIC EXTENSION of a kept-section label:
    there is a non-empty kept-section label P that is a proper prefix of the
    item's label and every component beyond P is a pure digit run ("A.1.1" =
    kept "A.1" + ["1"] -> promote; list marker "a" or mis-segmented prose
    "Fuehren"->[F,hren] -> NOT promoted). The anchors are the add-on's OWN kept
    section labels — nothing is hardcoded. Mutates the doc model in place (so body
    text is preserved for export) using the add-on's set_item_in_doc pattern, and
    returns the number of promotions."""
    from docling_core.types.doc.document import SectionHeaderItem, TextItem

    def label_of(item) -> tuple:
        return tuple(_segment_label(normalize_dashes((item.text or "").strip())))

    # Pass 1: trusted anchors = the add-on's kept section labels (non-empty).
    anchors: set[tuple] = set()
    for item, _ in doc.iterate_items(with_groups=False):
        if isinstance(item, SectionHeaderItem):
            lab = label_of(item)
            if lab:
                anchors.add(lab)

    # Pass 2: promote demoted TextItems whose label numerically extends an anchor.
    n_promo = 0
    for item, _ in list(doc.iterate_items(with_groups=False)):
        if isinstance(item, SectionHeaderItem) or not isinstance(item, TextItem):
            continue
        lab = label_of(item)
        if not lab:
            continue
        if _is_numeric_extension(lab, anchors):
            # TextItem -> SectionHeaderItem, then swap in at its self_ref index
            # (the add-on's set_item_in_doc pattern).
            header = SectionHeaderItem(
                **{
                    k: v
                    for k, v in item.model_dump().items()
                    if k != "label" and k in SectionHeaderItem.model_fields
                }
            )
            _, path, idx = item.self_ref.split("/")
            getattr(doc, path)[int(idx)] = header
            n_promo += 1
    return n_promo


def pdf_to_markdown(pdf_path: str) -> str:
    """Primary PDF route (INDEX-01-C1): pymupdf4llm -> relevel headings -> normalize dashes.
    Raises on empty/failed extraction so the caller can fall back to page_index (INDEX-01-C2)."""
    import pymupdf4llm

    # to_markdown() returns a str with default args; it only returns list[dict]
    # when page_chunks=True (which we do not pass). Cast to str for the type checker.
    md = cast(str, pymupdf4llm.to_markdown(pdf_path))
    if not md or not md.strip():
        raise RuntimeError(f"pdf_to_markdown produced empty output for {pdf_path}")
    return normalize_dashes(_relevel_headings(md))


class TessdataUnavailableError(RuntimeError):
    """Raised when non-Latin tessdata is missing and cannot be downloaded."""

    pass


_LATIN_LANGS = frozenset(
    {
        "afr",
        "cat",
        "ces",
        "dan",
        "deu",
        "eng",
        "est",
        "fin",
        "fra",
        "hrv",
        "hun",
        "ind",
        "isl",
        "ita",
        "lav",
        "lit",
        "msa",
        "nld",
        "nor",
        "pol",
        "por",
        "ron",
        "slk",
        "slv",
        "spa",
        "swe",
        "tur",
        "vie",
    }
)


# --- Fix 5: OCR language auto-detection + on-demand tessdata (RFC fizzy-forging-pearl) ---
# Deterministic, no model, no network for detection: classify by Unicode-script ratio.
_LATIN_LETTER_RE = re.compile(r"[A-Za-zÀ-ɏ]")
_DE_HINT_RE = re.compile(r"[äöüÄÖÜß]")
_AR_SCRIPT_MIN_RATIO = 0.15  # Arabic letters / all letters above which the doc is Arabic
_MIXED_SCRIPT_MIN_RATIO = 0.10  # Latin letters / all letters above which to add 'eng'
_AR_PRESENT_MIN_RATIO = 0.03  # any material Arabic -> include 'ara' (false-negative is costly)


def detect_ocr_langs(sample: str) -> list[str]:
    """Pick a Tesseract ``lang`` list from a text sample by Unicode-script ratio (Fix 5).

    Pure-Python, no dependency, no network (HR4). Returns Tesseract codes:
      * Arabic-dominant -> ['ara'] (or ['ara','eng'] when Latin is also materially
        present -- bilingual UAE gazettes);
      * German diacritics/ß present -> ['deu','eng'];
      * otherwise -> ['eng'].
    Empty / letterless input falls back to ['deu','eng'] to preserve the prior default.
    """
    text = sample or ""
    if not text.strip():
        return ["deu", "eng"]
    ar = len(_AR_SCRIPT_RE.findall(text))
    latin = len(_LATIN_LETTER_RE.findall(text))
    total = ar + latin
    if total == 0:
        return ["deu", "eng"]
    ar_ratio = ar / total
    if ar_ratio >= _AR_SCRIPT_MIN_RATIO:
        # Arabic-dominant: add 'eng' only when Latin is also materially present.
        return ["ara", "eng"] if latin / total >= _MIXED_SCRIPT_MIN_RATIO else ["ara"]
    if ar_ratio >= _AR_PRESENT_MIN_RATIO:
        # Latin-dominant but Arabic materially present (bilingual gazette) -> OCR both.
        return ["ara", "eng"]
    if _DE_HINT_RE.search(text):
        return ["deu", "eng"]
    return ["eng"]


def ensure_tessdata(langs: list[str]) -> list[str]:
    """Ensure ``<lang>.traineddata`` is available; return the usable subset (Fix 5).

    For each requested language, check ``TESSDATA_PREFIX`` for the traineddata file.
    Missing files are fetched from the official tessdata repo ONLY when
    ``TESSDATA_ALLOW_DOWNLOAD=1`` (egress-limited workers instead rely on PRE-BAKED
    traineddata in the image, mirroring the DOCLING_ARTIFACTS_PATH pre-bake). A
    missing Latin-script language is dropped (silent degrade is safe); a missing
    non-Latin-script language raises ``TessdataUnavailableError`` instead of being
    silently dropped, since that would silently degrade OCR to gibberish/empty
    output for scripts Latin OCR cannot read. If nothing remains after dropping
    Latin languages we fall back to ['deu','eng'] so OCR still runs. tessdata is
    data, not AGPL code (HR4)."""
    prefix = os.getenv("TESSDATA_PREFIX", "").strip()
    allow_dl = os.getenv("TESSDATA_ALLOW_DOWNLOAD", "0").strip().lower() in ("1", "true", "yes")
    available: list[str] = []
    for lang in langs:
        if not prefix:
            # No prefix configured -> trust the system tesseract install; assume present.
            available.append(lang)
            continue
        path = os.path.join(prefix, f"{lang}.traineddata")
        if os.path.exists(path):
            available.append(lang)
            continue
        if allow_dl and _try_download_tessdata(lang, prefix):
            available.append(lang)
        else:
            if lang not in _LATIN_LANGS:
                raise TessdataUnavailableError(
                    f"non-Latin tessdata missing: {lang} (prefix={prefix}, download={allow_dl})"
                )
            logger.warning(
                "tessdata for '%s' missing under %s (download=%s); dropping language",
                lang,
                prefix,
                allow_dl,
            )
    if not available:
        logger.warning("no requested OCR languages available; falling back to deu,eng")
        return ["deu", "eng"]
    return available


_TESSDATA_MAX_BYTES = 100 * 1024 * 1024  # 100 MB cap (RFC-009 D5 / Property 5)
_TESSDATA_CHUNK_BYTES = 1024 * 1024  # 1 MB chunked read
_TESSDATA_TIMEOUT_S = 30


def _try_download_tessdata(lang: str, prefix: str) -> bool:
    """Best-effort fetch of one traineddata file from the official repo. Never raises.

    Hardened per RFC-009 D5 (ISS-14): bounded by a 30s connection timeout and a
    100 MB total-size cap, both enforced via a chunked read loop. Any failure
    (timeout, oversize, network/HTTP error) cleans up the partial file at
    ``dest`` before returning False (Design Property 5: Tessdata download bounded).
    """
    import contextlib
    import urllib.request

    url = f"https://github.com/tesseract-ocr/tessdata/raw/main/{lang}.traineddata"
    dest = os.path.join(prefix, f"{lang}.traineddata")
    try:
        os.makedirs(prefix, exist_ok=True)
        total = 0
        with (
            urllib.request.urlopen(url, timeout=_TESSDATA_TIMEOUT_S) as resp,
            open(dest, "wb") as f,
        ):
            while True:
                chunk = resp.read(_TESSDATA_CHUNK_BYTES)
                if not chunk:
                    break
                total += len(chunk)
                if total > _TESSDATA_MAX_BYTES:
                    raise RuntimeError(
                        f"tessdata download for '{lang}' exceeded {_TESSDATA_MAX_BYTES} byte cap"
                    )
                f.write(chunk)
        logger.info("fetched tessdata for '%s' into %s (%d bytes)", lang, prefix, total)
        return True
    except Exception as exc:
        logger.warning("tessdata fetch failed for '%s' (%s)", lang, exc)
        if os.path.exists(dest):
            with contextlib.suppress(OSError):
                os.unlink(dest)
        return False


def _build_pdf_pipeline_options(
    force_full_page_ocr: bool = False,
    ocr_lang_override: list[str] | None = None,
):
    """Build the CPU-only Docling PDF pipeline options.

    Fix 3 (RFC fizzy-forging-pearl): ``force_full_page_ocr`` re-OCRs the WHOLE page
    even when a (corrupt) text layer is present -- the only way Docling will overwrite
    a broken CMap/font text layer (the مرسوم class). ``ocr_lang_override`` pins the
    Tesseract language list (from Fix-5 detection) instead of the static env default.
    Both default to the prior behaviour, so the normal converter path is unchanged.
    ``DOCLING_FORCE_FULL_PAGE_OCR=1`` is honoured as a manual override.

    Capping intra-op threads (``DOCLING_NUM_THREADS``, default 1) is the one
    code-level RSS reducer that costs NO extraction fidelity: Docling propagates
    ``num_threads`` to ``torch.set_num_threads`` / onnxruntime internally, so peak
    memory drops (fewer per-thread scratch arenas) without unloading any model or
    changing output. TableFormer stays on at ``ACCURATE`` -- disabling it or using
    ``FAST`` would cut memory further but degrade table reconstruction, which we do
    NOT want. Docling imports stay function-local (they are heavy).
    """
    from docling.datamodel.accelerator_options import AcceleratorDevice, AcceleratorOptions
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
        TableFormerMode,
        TesseractCliOcrOptions,
    )

    # CPU-only by design -- nothing on GPU/MPS for now.
    device = AcceleratorDevice.CPU
    # Fix 3: full-page OCR (param or DOCLING_FORCE_FULL_PAGE_OCR=1) forces do_ocr on so a
    # corrupt existing text layer can be overwritten; it implies do_ocr regardless of env.
    force_ocr = force_full_page_ocr or os.getenv(
        "DOCLING_FORCE_FULL_PAGE_OCR", "0"
    ).strip().lower() in ("1", "true", "yes")
    do_ocr = force_ocr or os.getenv("DOCLING_DO_OCR", "0").strip().lower() in ("1", "true", "yes")
    # Cap inference threads to bound peak RSS. Default 1 for the memory-tight worker;
    # raise via DOCLING_NUM_THREADS only where the node has RAM headroom.
    try:
        num_threads = max(1, int(os.getenv("DOCLING_NUM_THREADS", "1")))
    except ValueError:
        num_threads = 1

    opts = PdfPipelineOptions()
    opts.do_ocr = do_ocr
    opts.do_table_structure = True
    opts.table_structure_options.mode = TableFormerMode.ACCURATE
    if do_ocr:
        # Fix 5: an explicit detected-language override beats the static env list.
        langs = ocr_lang_override or [
            s.strip() for s in os.getenv("DOCLING_OCR_LANG", "deu,eng").split(",") if s.strip()
        ]
        # CLI engine -> uses the system `tesseract` binary, which honours TESSDATA_PREFIX.
        opts.ocr_options = TesseractCliOcrOptions(lang=langs, force_full_page_ocr=force_ocr)
    opts.accelerator_options = AcceleratorOptions(device=device, num_threads=num_threads)
    # Use pre-baked model artifacts when available (set in the container image so
    # egress-limited workers never download weights at runtime -- a download failure
    # there would otherwise raise and silently fall back to pymupdf4llm -> flat tree
    # -> depth<2). Unset (local dev) -> docling fetches from HF on first use.
    artifacts_path = os.getenv("DOCLING_ARTIFACTS_PATH", "").strip()
    if artifacts_path:
        opts.artifacts_path = artifacts_path
    return opts


# Process-lifetime DocumentConverter cache. Constructing a DocumentConverter
# loads the Heron layout + TableFormer model weights (~700 MB-1.4 GB RSS) and
# torch NEVER returns that to the OS. Building a NEW one per pdf_to_markdown_docling
# call therefore leaks ~250 MB per document in any process that converts more than
# once in-process (e.g. preprocess_client.py, which asyncio.gathers the whole
# doc_store) — RSS climbs monotonically until OOM. Reusing one instance caps growth
# to a few MB/doc (measured 2026-06-13: +237 MB/call rebuilt vs +5-7 MB/call reused).
# The production worker is unaffected — each job runs in a fresh converters_cli child
# that dies — but this keeps the in-process callers bounded. Keyed on the env knobs
# _build_pdf_pipeline_options() reads so a mid-process env change rebuilds correctly.
_DOCLING_CONVERTER_CACHE: dict[tuple[str, ...], "DocumentConverter"] = {}


def _docling_converter(
    force_full_page_ocr: bool = False,
    ocr_lang_override: list[str] | None = None,
    for_image: bool = False,
) -> "DocumentConverter":
    """Return a cached CPU-only DocumentConverter, building it once per options key.

    Docling's DocumentConverter is designed for reuse across .convert() calls: the
    models load on first use and are reused, so a single instance is both correct
    and the memory-bounded choice. See _DOCLING_CONVERTER_CACHE above for why a new
    instance per call leaks.

    Fix 3: the optional force-OCR / language-override flags are part of the cache key
    so an escalation converter is a distinct, separately-cached instance and the normal
    (no-arg) path keeps its existing key and cached object untouched.

    Fix 5: ``for_image`` routes InputFormat.IMAGE instead of InputFormat.PDF through the
    same StandardPdfPipeline options and is part of the cache key, so image_to_markdown()
    shares this process-lifetime cache instead of building a fresh (leaking) converter.
    """
    from docling.datamodel.base_models import InputFormat
    from docling.document_converter import DocumentConverter, PdfFormatOption

    key = (
        os.getenv("DOCLING_DO_OCR", "0").strip().lower(),
        os.getenv("DOCLING_NUM_THREADS", "1").strip(),
        os.getenv("DOCLING_OCR_LANG", "deu,eng").strip(),
        os.getenv("DOCLING_ARTIFACTS_PATH", "").strip(),
        "force" if force_full_page_ocr else "",
        ",".join(ocr_lang_override) if ocr_lang_override else "",
        "image" if for_image else "pdf",
    )
    converter = _DOCLING_CONVERTER_CACHE.get(key)
    if converter is None:
        pipeline_options = _build_pdf_pipeline_options(
            force_full_page_ocr=force_full_page_ocr,
            ocr_lang_override=ocr_lang_override,
        )
        input_format = InputFormat.IMAGE if for_image else InputFormat.PDF
        converter = DocumentConverter(
            format_options={input_format: PdfFormatOption(pipeline_options=pipeline_options)}
        )
        _DOCLING_CONVERTER_CACHE[key] = converter
        logger.info("instantiated and cached Docling DocumentConverter (options key=%s)", key)
    return converter


_HIERARCHICAL_INFER_PATCHED = False
# Fingerprint of the upstream strict-equality match we replace. If the installed
# docling-hierarchical-pdf version no longer contains this line, we skip the patch
# rather than risk a stale override — the Rank-1 over-prune fallback covers us.
_HBM_STRICT_MATCH_FINGERPRINT = 're.sub(r"[^A-Za-z0-9]", "", title) == re.sub('
# A TOC title must have at least this many alphanumerics before we accept a
# numbering-prefix suffix match, to avoid short-word false positives
# (e.g. bare "Tierhaltung" suffix-matching an unrelated heading).
_HBM_MIN_SUFFIX_LEN = 5


# Complexity grandfathered (hierarchical add-on patch); see pyproject [tool.ruff].
def _patch_hierarchical_infer() -> None:  # noqa: C901, PLR0915
    """Make the docling-hierarchical-pdf add-on tolerate publisher numbering prefixes.

    The add-on's ``HierarchyBuilderMetadata.infer()`` matches PDF-outline (TOC)
    titles to Docling document items by STRICT stripped-alphanumeric equality
    (hierarchy_builder_metadata.py:189). German insurance PDFs (e.g. the BHB
    Haftpflicht booklet) list bare titles in the TOC ("Land- und Forstwirtschaft")
    while the in-document heading carries a clause prefix ("BHB 3 Land- und
    Forstwirtschaft"), so the equality fails for ~32/33 entries and the add-on
    demotes almost every heading to body text -> node_count<3 rejection (HR5).

    This installs (once, idempotently) a patched ``infer()`` whose matching falls
    back to a SUFFIX match (``item.orig`` ends with the TOC title) when no exact
    match exists, guarded by ``_HBM_MIN_SUFFIX_LEN`` and constrained to the TOC
    entry's target page (the loop already iterates ``page_no=page``); among suffix
    candidates it prefers the shortest item (least extra prefix) so a real heading
    "BHB 3 X" wins over a longer body sentence ending in "X". The patch is
    fingerprint-guarded against upstream version drift, and the caller wraps it so
    it can NEVER be fatal — on any failure the gate-aware source selection in
    ``pdf_to_markdown_docling`` falls back to raw Docling markdown.
    """
    global _HIERARCHICAL_INFER_PATCHED
    if _HIERARCHICAL_INFER_PATCHED:
        return
    import inspect

    from docling_core.types.doc.document import ListItem, TextItem
    from hierarchical import hierarchy_builder_metadata as _hbm
    from hierarchical.hierarchy_builder_metadata import (
        HeaderNotFoundException,
        HierarchyBuilderMetadata,
        ImplausibleHeadingStructureException,
    )
    from hierarchical.types.hierarchical_header import HierarchicalHeader

    src = inspect.getsource(HierarchyBuilderMetadata.infer)
    if _HBM_STRICT_MATCH_FINGERPRINT not in src:
        logger.warning(
            "hierarchical infer() match logic changed upstream; skipping "
            "suffix-match patch (relying on raw-docling over-prune fallback)"
        )
        _HIERARCHICAL_INFER_PATCHED = True  # don't re-inspect on every conversion
        return

    def _patched_infer(self) -> HierarchicalHeader:
        # Copy of HierarchyBuilderMetadata.infer with the item-matching loop made
        # tolerant of a missing numbering prefix; the rest is upstream-verbatim.
        heading_to_level = self._extract_toc()
        root = HierarchicalHeader()
        current = root
        doc = self.conv_res.document

        for level, title, page, add_info in heading_to_level:
            new_parent = None
            this_item = None
            title_norm = re.sub(r"[^A-Za-z0-9]", "", title)
            suffix_item = None
            suffix_norm_len: int | None = None
            for item, _ in doc.iterate_items(page_no=page):
                if not isinstance(item, (TextItem, ListItem)):
                    continue
                item_norm = re.sub(r"[^A-Za-z0-9]", "", item.orig)
                if item_norm == title_norm:
                    this_item = item  # exact match always wins
                    break
                # numbering-prefix-tolerant fallback: keep the tightest suffix match
                if (
                    len(title_norm) >= _HBM_MIN_SUFFIX_LEN
                    and item_norm.endswith(title_norm)
                    and (suffix_norm_len is None or len(item_norm) < suffix_norm_len)
                ):
                    suffix_item = item
                    suffix_norm_len = len(item_norm)
            if this_item is None:
                this_item = suffix_item
            if this_item is None:
                if self.raise_on_error:
                    raise HeaderNotFoundException(add_info)
                else:
                    _hbm.logger.warning(HeaderNotFoundException(add_info))
                    continue

            if current.level_toc is None or level > current.level_toc:
                new_parent = current
            elif level == current.level_toc:
                if current.parent is not None:
                    new_parent = current.parent
                else:
                    raise ImplausibleHeadingStructureException()
            else:
                new_parent = current
                while new_parent.parent is not None and (level <= new_parent.level_toc):
                    new_parent = new_parent.parent
            new_obj = HierarchicalHeader(
                text=this_item.orig,
                parent=new_parent,
                level_toc=level,
                doc_ref=this_item.self_ref,
            )
            new_parent.children.append(new_obj)
            current = new_obj

        return root

    HierarchyBuilderMetadata.infer = _patched_infer
    _HIERARCHICAL_INFER_PATCHED = True
    logger.info(
        "patched hierarchical infer() with numbering-prefix suffix matching (min title len %d)",
        _HBM_MIN_SUFFIX_LEN,
    )


_INDENTED_HEADING_RE = re.compile(r"^[ \t]+(#{1,6}\s)", re.MULTILINE)


def _normalize_indented_headings(md: str) -> str:
    """Strip leading whitespace before markdown heading markers."""
    return _INDENTED_HEADING_RE.sub(r"\1", md)


# RFC-015 D5c: Docling occasionally emits several '#{1,6} '-prefixed headings on ONE
# physical line ("text### Heading"); the line-anchored heading regexes downstream see
# only the first, so the rest fuse into a giant tail blob (doc 7dcf7cb7). The lookbehind
# `[^\n#]` matches only a heading marker preceded by a non-newline, non-'#' char, so a
# marker mid-line ("text### X") splits while true line-start headings AND the interior
# of a genuine multi-'#' run ("## X" must NOT split into "#\n# X") are left untouched.
# (Tightened from the RFC-015 pseudocode's `[^\n]`, which wrongly split multi-'#' runs.)
_RUN_TOGETHER_HEADING_RE = re.compile(r"(?<=[^\n#])(#{1,6}\s)")


def _split_run_together_headings(md: str) -> str:
    """Insert a newline before any heading marker Docling emitted mid-line (RFC-015 D5c).

    Runs BEFORE the hash-sentinel fix (D4) and heading-depth inference so a run-together
    ``text### Heading`` is separated onto its own line first — both D4's per-line pass
    and the splitter's ordinal matching only inspect one marker per line."""
    return _RUN_TOGETHER_HEADING_RE.sub(r"\n\1", md)


# RFC-015 D4: consume WHOLE '#+' runs (not just interior '#'). RFC-010 D5 used
# `(?<=\S)#(?=\S)` (non-whitespace both sides), so when the corrupted في run's outer
# edges sat next to whitespace — the normal case, في being a standalone word — the
# boundary '#'s survived, leaving `#في#`/`#فيفي#`. That residue poisoned the splitter's
# oversized-ordinal anchor, fusing whole legal instruments into one leaf (doc aebf15b4).
_INLINE_HASH_RE = re.compile(r"#+")
_HEADING_MARKER_LINE_RE = re.compile(r"#{1,6}[ \t]")


def _fix_fi_hash_substitution(md: str) -> str:
    """Restore في from Docling's ``#`` substitution in Arabic text (RFC-015 D4).

    Docling renders the standalone Arabic word في as ``#``. This converts every inline
    ``#+`` run back to في, per line, while preserving genuine line-initial markdown
    heading markers (``#{1,6}`` followed by a space). Widened from the RFC-010 D5
    interior-only regex so boundary/standalone ``#`` are also consumed (see the note
    above). Runs EARLIER in the pipeline (before heading-depth inference) so في is a
    single token by the time the heading regex sees it. Pure local string surgery — no
    LLM, no network (HR3)."""
    if not md:
        return md
    arabic = sum(1 for c in md if "؀" <= c <= "ۿ")
    if arabic / len(md) <= 0.15:
        return md
    out: list[str] = []
    for line in md.splitlines(keepends=True):
        stripped = line.lstrip()
        # Preserve a genuine markdown heading marker ("## Heading"); convert everything
        # else — inline '#' runs that are corrupted في — to في.
        if _HEADING_MARKER_LINE_RE.match(stripped):
            out.append(line)
        else:
            out.append(_INLINE_HASH_RE.sub("في", line))
    return "".join(out)




def reconstruct_bidi_order(text: str, expected_script: str | None = None) -> str:
    """Zone-3: apply_rtl shim replacing the old bidi reconstructor.

    Two-level strategy (preserves RFC-023 D9 per-heading correction):
    1. Document-level: if decide_rtl says the whole text is reversed,
       apply_rtl repairs all lines.
    2. Per-heading: even when the document is NOT reversed overall,
       each heading line is checked individually — a visual-order
       heading in an otherwise-logical document still gets corrected.

    ``expected_script`` is accepted for call-site compatibility but is
    unused (``decide_rtl`` infers script from content).
    """
    if not text:
        return text
    arabic = len(_AR_SCRIPT_RE.findall(text))
    if arabic == 0:
        return text

    decision: RtlDecision = decide_rtl(text)
    if decision.reversed:
        return apply_rtl(text, reversed_flag=True)

    out: list[str] = []
    changed = False
    for line in text.splitlines(keepends=True):
        m = _BIDI_HEADING_PREFIX_RE.match(line)
        if m:
            heading_text = m.group(2)
            if decide_rtl(heading_text.strip(), sample_count=1).reversed:
                repaired = apply_rtl(heading_text.rstrip(), reversed_flag=True)
                eol = line[len(line.rstrip()):]
                out.append(m.group(1) + repaired + eol)
                changed = True
                continue
        out.append(line)
    return "".join(out) if changed else text


# Split a leading markdown heading marker off a line so reconstruct_bidi_order reorders
# only the title text, leaving the '#' prefix in place for depth inference.
_BIDI_HEADING_PREFIX_RE = re.compile(r"^(\s*#{1,6}[ \t]+)(.*)$", re.DOTALL)



# Zone-3: _fix_residual_rtl_reversal DELETED — redundant with
# reconstruct_bidi_order which already covers decide_rtl+apply_rtl.


# RFC-015 D6 gate. Consolidated in config.py (canonical source); imported here
# to eliminate the prior client.py / converters.py double-definition.
# Zone-5: per-picture enrichment uses the dedicated OCR_ESCALATION_PER_PICTURE flag.
from .config import OCR_ESCALATION_PER_PICTURE as _OCR_ESCALATION_PER_PICTURE

_IMAGE_MARKER = "<!-- image -->"
_PICTURE_OCR_MIN_CHARS = 20  # RFC-015 D6: below this, OCR output is decorative-image noise
_PICTURE_PAGE_COVERAGE_THRESHOLD = float(os.getenv("PICTURE_PAGE_COVERAGE_THRESHOLD", "0.6"))
# D2 (RFC-023): sub-icon PictureItems (both dims below this) skip crop+OCR
# entirely and are tagged "decorative_icon" — set to 0 to disable the pre-filter.
_DECORATIVE_ICON_MIN_DIM_PT = float(os.getenv("DECORATIVE_ICON_MIN_DIM_PT", "20"))
# Audit 2026-07-21 finding 10: bound for the per-picture OCR and VLM thread pools.
# Keeps a many-figure document from spawning unbounded tesseract subprocesses or
# parallel paid vision calls inside one conversion.
_IMAGE_ENRICH_CONCURRENCY = max(1, int(os.getenv("IMAGE_ENRICH_CONCURRENCY", "4") or "4"))
# F1 (RFC-020): when True, pages with no text layer are exempt from the coverage
# skip — the full-page picture IS the content and must be OCR'd.
_COVERAGE_EXEMPT_NO_TEXT_LAYER = os.getenv(
    "COVERAGE_EXEMPT_NO_TEXT_LAYER", "true"
).strip().lower() in ("1", "true", "yes")
# D0 (RFC-023): when True, a text layer that passes the char-count check but is
# garbled (mojibake/scanned-PDF noise) is still treated as "no content", so the
# coverage exemption above fires and per-picture OCR proceeds.
_TEXT_LAYER_GARBLE_CHECK_ENABLED = os.getenv(
    "TEXT_LAYER_GARBLE_CHECK_ENABLED", "true"
).strip().lower() in ("1", "true", "yes")
# D1 (RFC-024): capture clip_text into PictureResult.ocr_text when it is NOT
# already contained in the Docling markdown export (containment-guarded — see
# _clip_text_contained). Set to False to restore the pre-RFC-024 skip-only
# behavior (every non-trivial clip_text is discarded, "clip_text" reason).
_CLIP_TEXT_CAPTURE_ENABLED = os.getenv("CLIP_TEXT_CAPTURE_ENABLED", "true").strip().lower() in (
    "1",
    "true",
    "yes",
)
_CLIP_TEXT_CONTAINMENT_THRESHOLD = 0.6
# D1 (RFC-024): document-level text-layer fallback — when Docling's exported
# markdown is this thin (excluding `<!-- image -->` markers), the document is
# image-dominant and the native PDF text layer is read wholesale as
# supplementary content rather than lost entirely.
_DOC_TEXT_FALLBACK_MIN_CHARS = 100
# D1 (RFC-025): secondary trigger — a heading-only tree can carry enough total
# chars to clear _DOC_TEXT_FALLBACK_MIN_CHARS while every heading has almost no
# body prose beneath it (structure survived, content did not). Fire the same
# pdfium whole-document fallback when chars-per-heading drops below this floor.
_DOC_TEXT_FALLBACK_MIN_CHARS_PER_HEADING = 50
# D1 (RFC-025): when True, the full-page-coverage exemption uses the
# region-scoped text check (_region_has_own_text_layer); when False, restore
# the pre-RFC-025 page-level check (_text_layer_has_content) for rollback.
_REGION_AWARE_TEXT_CHECK_ENABLED = os.getenv(
    "REGION_AWARE_TEXT_CHECK_ENABLED", "true"
).strip().lower() in ("1", "true", "yes")
# D1 (RFC-025): the region-aware exemption converts previously-skipped
# full-page picture regions into active 300-DPI crop+Tesseract OCR work,
# which is expensive on multi-hundred-page scanned documents. Cap the number
# of full-page exemptions fired per document; further regions past the cap
# are skipped (page_coverage) with a warning.
_MAX_FULLPAGE_PICTURE_OCR_REGIONS = int(os.getenv("MAX_FULLPAGE_PICTURE_OCR_REGIONS", "50"))


class PictureResult(TypedDict, total=False):
    """Structured result from per-picture OCR/crop recovery."""

    ocr_text: str
    png_bytes: bytes
    page: int
    bbox: dict
    description: str
    skipped_reason: str  # RFC-019 D3: deliberate-skip tag (e.g. "page_coverage")
    decorative: bool  # RFC-023 D2: sub-icon pre-filter or empty-OCR belt-and-suspenders


def zdr_egress_gate(purpose: str, doc_id: str = "") -> tuple[bool, str | None]:
    """Shared HR3 gate for every image/doc-text LLM egress (audit findings 2/3).

    Returns ``(allowed, api_base)``. ``api_base`` is the SAME endpoint the caller
    MUST pass to ``litellm.completion(api_base=...)`` so the gate inspects exactly
    what egresses — litellm resolving a different endpoint from its own env would
    otherwise silently diverge from the inspected one (finding 3). Blocks when
    ``pii_corpus`` is set and the endpoint is not on the ZDR allow-list."""
    from .config import _is_zdr_allowlisted, settings

    api_base = settings.openai_base_url
    if settings.pii_corpus and not _is_zdr_allowlisted(api_base):
        logger.info(
            "%s skipped for %s: pii_corpus=True, endpoint not ZDR-allowlisted (HR3)",
            purpose,
            doc_id or "<unknown doc>",
        )
        return False, api_base
    return True, api_base


def _collect_picture_regions(doc) -> list[dict]:
    """List each PictureItem's 1-indexed page + bbox in document iteration order (D6).

    The order matches the ``<!-- image -->`` markers ``export_to_markdown()`` emits, so
    the caller can splice recovered text by positional index (picture bboxes are stable
    across the add-on's in-place mutation, unlike heading selection)."""
    from docling_core.types.doc.document import PictureItem

    regions: list[dict] = []
    for item, _ in doc.iterate_items(with_groups=False):
        if isinstance(item, PictureItem) and item.prov:
            prov = item.prov[0]
            regions.append({"page": prov.page_no, "bbox": prov.bbox})
    return regions


def _bbox_to_fitz_rect(bbox, page_height: float, fitz):
    """Convert a Docling BoundingBox to a top-left-origin ``fitz.Rect`` (D6).

    Docling bboxes may carry a BOTTOMLEFT coordinate origin (PDF-native), while
    ``fitz.Rect`` is TOP-LEFT; convert using the page height when needed. Returns None
    on any unusable bbox so the caller skips that picture."""
    try:
        left, top, right, bottom = bbox.l, bbox.t, bbox.r, bbox.b
        origin = getattr(bbox, "coord_origin", None)
        origin_name = getattr(origin, "name", str(origin or "")).upper()
        if origin_name.startswith("BOTTOM"):
            top, bottom = page_height - top, page_height - bottom
        y0, y1 = sorted((top, bottom))
        x0, x1 = sorted((left, right))
        if x1 - x0 <= 0 or y1 - y0 <= 0:
            return None
        return fitz.Rect(x0, y0, x1, y1)
    except Exception:
        return None


def _tesseract_ocr_image(png_path: str, langs: list[str]) -> str:
    """OCR one image file via the LOCAL ``tesseract`` CLI (RFC-015 D6; HR3-clean).

    Uses the same system ``tesseract`` binary + ``TESSDATA_PREFIX`` the Docling OCR
    path uses — no LLM, no network egress, so PII in a chart never leaves the host
    (HR3). Returns stripped recognised text, or '' on any failure (never raises)."""
    tess = shutil.which("tesseract")
    if not tess:
        logger.warning("tesseract binary not found; skipping per-picture OCR")
        return ""
    try:
        proc = subprocess.run(
            [tess, png_path, "stdout", "-l", "+".join(langs)],
            capture_output=True,
            text=True,
            timeout=60,
        )
        return proc.stdout.strip()
    except Exception as exc:
        logger.warning("per-picture tesseract OCR failed (%s)", exc)
        return ""


def _text_layer_has_content(page, expected_script: str | None = None) -> bool:
    """Return True when the page's native text layer has meaningful content.

    Used by F1 (RFC-020) to exempt scanned pages from the coverage skip —
    when a page has NO text layer the full-page picture IS the content.

    D0 (RFC-023): a text layer can pass the char-count check but still be
    garbled (thin mojibake left by the PDF creator on a scanned page). Treat
    a garbled text layer as no content so the coverage exemption fires."""
    text = page.get_text("text").strip()
    if len(text) <= _PICTURE_OCR_MIN_CHARS:
        return False
    if _TEXT_LAYER_GARBLE_CHECK_ENABLED:
        from .helpers import GarbleContext, check_garble
        from .script import infer_script

        if check_garble(text, expected_script=expected_script or infer_script(text), context=GarbleContext.PAGE_TEXT_LAYER):
            return False
    return True


def _region_has_own_text_layer(page, region_rect) -> bool:
    """Return True when the picture region's OWN bbox has a meaningful native text layer.

    D1 (RFC-025): `_text_layer_has_content(page)` is a page-level check, so
    incidental text outside the region (headers, footers, page numbers) keeps
    it True and disables the coverage exemption even when the region's own
    body text was baked into a skipped full-page image. This checks only the
    text clipped to `region_rect`, regardless of what text exists outside it."""
    region_clip_len = len(page.get_text("text", clip=region_rect).strip())
    return region_clip_len >= _PICTURE_OCR_MIN_CHARS


def _normalize_for_containment(text: str) -> str:
    """NFKC-fold + whitespace-collapse + lowercase (RFC-024 D1).

    Shared between the clip-text containment guard and its tests so both sides
    of the comparison are robust to whitespace/reflow differences between the
    PDF text layer and the Docling markdown export."""
    return " ".join(unicodedata.normalize("NFKC", text).split()).lower()


def _clip_text_contained(clip_text: str, md_norm: str) -> bool:
    """True when >=60% of ``clip_text``'s normalized chars belong to tokens
    that appear as substrings of ``md_norm`` (RFC-024 D1 containment guard
    against double-capturing content Docling already exported into the
    markdown body). Token-level substring matching (length-weighted) stays
    robust to whitespace/reflow differences while remaining discriminative —
    a raw character-frequency check would report near-total containment for
    any clip against a large markdown body."""
    clip_norm = _normalize_for_containment(clip_text)
    if not clip_norm:
        return True
    if not md_norm:
        return False
    tokens = clip_norm.split()
    total = sum(len(t) for t in tokens)
    if total == 0:
        return True
    matched = sum(len(t) for t in tokens if t in md_norm)
    return matched / total >= _CLIP_TEXT_CONTAINMENT_THRESHOLD


def _document_level_text_fallback(md: str, pdf_path: str, expected_script: str | None = None) -> str:
    """Full-page text-layer fallback for image-dominant documents (RFC-024 D1).

    When Docling's exported markdown carries fewer than
    ``_DOC_TEXT_FALLBACK_MIN_CHARS`` characters (excluding ``<!-- image -->``
    markers), Docling routed nearly the entire document through the picture
    path and per-region recovery in ``_recover_picture_text`` has no markdown
    body to work against. Read the native PDF text layer wholesale via
    pypdfium2 (BSD-3/Apache-2, HR4) and append it as supplementary content so
    the tree build sees something other than bare image markers. Never fatal
    — any failure returns ``md`` unchanged.

    Secondary trigger (RFC-025 D1): a heading-only tree (structure survived,
    body prose did not — e.g. a 347-node ToC with no article text) can clear
    the total-char floor above while still carrying almost no prose per
    heading. Fire the same fallback when chars-per-heading drops below
    ``_DOC_TEXT_FALLBACK_MIN_CHARS_PER_HEADING``."""
    total_chars = len(md.replace(_IMAGE_MARKER, ""))
    heading_count = len(_HEADING_RE.findall(md))
    if (
        total_chars >= _DOC_TEXT_FALLBACK_MIN_CHARS
        and total_chars / max(heading_count, 1) >= _DOC_TEXT_FALLBACK_MIN_CHARS_PER_HEADING
    ):
        return md
    try:
        import pypdfium2 as pdfium

        pdoc = pdfium.PdfDocument(pdf_path)
        try:
            page_texts = []
            for page in pdoc:
                textpage = page.get_textpage()
                text = textpage.get_text_range().strip()
                if text:
                    page_texts.append(text)
        finally:
            pdoc.close()
    except Exception as exc:
        logger.warning(
            "document-level text-layer fallback failed for %s (%s); keeping markdown as-is",
            pdf_path,
            exc,
        )
        return md
    full_text = "\n\n".join(page_texts).strip()
    if not full_text:
        return md
    # RFC-024 D1 risk mitigation: a scanned page can carry a thin mojibake text
    # layer — never append a garbled text layer as supplementary content (HR5).
    from .helpers import GarbleContext, check_garble
    from .script import infer_script

    if check_garble(full_text, expected_script=expected_script or infer_script(full_text), context=GarbleContext.DOCUMENT_FALLBACK):
        logger.warning(
            "document-level text-layer fallback skipped for %s: text layer is garbled",
            pdf_path,
        )
        return md
    logger.info(
        "document-level text-layer fallback fired for %s (%d markdown char(s) "
        "excluding image markers)",
        pdf_path,
        total_chars,
    )
    return f"{md}\n\n{full_text}"


def _page_rotation_correction_info(page) -> dict:
    """RFC-026 D2: read a single page's /Rotate metadata plus an aspect-ratio
    fallback. Reuses the `page.rotation` accessor already used at the D6
    crop-normalization site (~line 1746). Per-page, not per-document — a
    single PDF can mix portrait and landscape pages. `/Rotate` is authoritative;
    the aspect-ratio heuristic is advisory only, for pages where a scanner
    omitted `/Rotate` but still produced a wide page.
    """
    try:
        rotate = page.rotation
        width = page.rect.width
        height = page.rect.height
    except Exception:
        return {"rotate": 0, "likely_landscape": False, "width": 0.0, "height": 0.0}
    likely_landscape = rotate == 0 and width > height
    return {
        "rotate": rotate,
        "likely_landscape": likely_landscape,
        "width": width,
        "height": height,
    }


# RFC-026 D2: gates the rotation-aware coordinate transform below so the fix
# can be rolled back without a revert if it regresses an unrelated corpus doc.
_PAGE_ROTATION_DETECTION_ENABLED = os.getenv(
    "PAGE_ROTATION_DETECTION_ENABLED", "true"
).strip().lower() in ("1", "true", "yes")


def _normalize_pdf_page_rotation(pdf_path: str) -> str:
    """RFC-026 D2: bake each page's effective rotation into its `/Rotate` key
    before handing the file to the docling extraction backend, so rotated and
    aspect-ratio-landscape pages get a consistent coordinate mapping instead of
    fragmenting into near-empty text blocks (the `uae_numbers_english_page_16_17`
    stall — ~750 chars extracted vs. ~4000-8000 expected).

    `effective_rotation` is `/Rotate` when non-zero, else 90 when the aspect-ratio
    heuristic flags a likely-landscape page with no explicit `/Rotate` (a scanner
    that omitted the key but still produced a wide page). Per-page — mixed
    portrait/landscape documents only get pages that need it rewritten.

    Composes with the existing D6 rotation-zeroing at the OCR-crop site
    (~lines 1744-1774): that path opens its own `fitz.Document` for cropping and
    always restores `orig_rotation` after rendering, so it is unaffected by the
    (separate, disk-persisted) copy this function may return.

    Returns the original path unchanged when no page needs correction, when the
    gate is disabled, or on any read/write failure (fail-open — a single
    corrupted page's rotation metadata must not abort extraction).
    """
    if not _PAGE_ROTATION_DETECTION_ENABLED:
        return pdf_path
    from .config import ALLOW_AGPL_FALLBACK

    if not ALLOW_AGPL_FALLBACK:
        logger.warning(
            "rotation normalization skipped for %s: ALLOW_AGPL_FALLBACK=false "
            "(fitz/PyMuPDF is AGPL-3.0)",
            pdf_path,
        )
        return pdf_path
    try:
        import fitz  # PyMuPDF, AGPL-3.0

        pdf = fitz.open(pdf_path)
        try:
            changed = False
            for page in pdf:
                info = _page_rotation_correction_info(page)
                effective_rotation = (
                    info["rotate"] if info["rotate"] else (90 if info["likely_landscape"] else 0)
                )
                if effective_rotation and effective_rotation != page.rotation:
                    page.set_rotation(effective_rotation)
                    changed = True
            if not changed:
                return pdf_path
            # SIM115 rationale: the temp FILE must outlive this scope -- its path is
            # returned to the caller, who reads it and unlinks it later. A context
            # manager would delete/close it before the caller ever opens it.
            tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)  # noqa: SIM115
            tmp.close()
            pdf.save(tmp.name)
            return tmp.name
        finally:
            pdf.close()
    except Exception as exc:
        logger.warning(
            "rotation normalization failed for %s (%s); using original file", pdf_path, exc
        )
        return pdf_path


def _tag_landscape_pages_for_fallback(pdf_path: str) -> list[dict]:
    """RFC-035 D2 Phase 1: read-only pre-extraction landscape orientation probe.

    Tags each page landscape via PyMuPDF's `page.rotation` (rotation % 180 != 0)
    OR a `width > height` geometric heuristic, threading the result as
    extraction metadata for the Phase 2 rasterize-rotate-reextract fallback.
    Does not mutate the PDF or the primary extraction path — portrait-only
    documents yield an all-False probe and are otherwise unaffected.
    """
    from .config import ALLOW_AGPL_FALLBACK

    if not ALLOW_AGPL_FALLBACK:
        return []
    try:
        import fitz  # PyMuPDF, AGPL-3.0

        pages = []
        with fitz.open(pdf_path) as doc:
            for page_no, page in enumerate(doc):
                try:
                    rotate = page.rotation
                    width = page.rect.width
                    height = page.rect.height
                    is_landscape = (rotate % 180 != 0) or (width > height)
                except Exception:
                    rotate, width, height, is_landscape = 0, 0.0, 0.0, False
                pages.append(
                    {
                        "page_no": page_no,
                        "rotate": rotate,
                        "width": width,
                        "height": height,
                        "is_landscape": is_landscape,
                    }
                )
        return pages
    except Exception as exc:
        logger.warning("landscape orientation probe failed for %s (%s)", pdf_path, exc)
        return []


# RFC-035 D2 Phase 2 trigger: below this char count, a landscape-tagged page's
# primary extraction is considered failed and the rasterize-rotate-reextract
# fallback should engage. Configurable per corpus (chart-heavy pages may need
# a different floor than the 748-char stalled baseline this threshold targets).
LANDSCAPE_CHAR_THRESHOLD: int = int(os.environ.get("LANDSCAPE_CHAR_THRESHOLD", "500"))

# RFC-036 D0a: hard caps on the per-page reextraction loop below, so a
# document with many low-char landscape pages cannot serially rasterize/OCR
# its way past the chunk timeout budget.
MAX_LANDSCAPE_PAGES: int = int(os.environ.get("MAX_LANDSCAPE_PAGES", "10"))
LANDSCAPE_REEXTRACT_DEADLINE_SECONDS: float = float(
    os.environ.get("LANDSCAPE_REEXTRACT_DEADLINE_SECONDS", "600")
)


def _landscape_pages_below_threshold(document, landscape_pages: list[dict]) -> list[dict]:
    """RFC-035 D2 Phase 2 trigger: for pages tagged landscape by
    ``_tag_landscape_pages_for_fallback``, count the chars Docling's primary extraction
    yielded for that page and flag pages below ``LANDSCAPE_CHAR_THRESHOLD``
    that also have a detectable picture/graphic region (RFC-036 D0c) as
    needing the rasterize-rotate-reextract fallback. Dense numeric-table
    pages (e.g. world-stats-pocketbook) fall below the char threshold but
    carry no picture region, so they no longer false-positive trigger.
    """
    if not any(p["is_landscape"] for p in landscape_pages):
        return []
    picture_pages = {r["page"] for r in _collect_picture_regions(document)}
    below = []
    for p in landscape_pages:
        if not p["is_landscape"]:
            continue
        # PyMuPDF page_no is 0-indexed; Docling's prov.page_no (and
        # iterate_items' page_no kwarg) is 1-indexed.
        page_no = p["page_no"] + 1
        if page_no not in picture_pages:
            continue
        char_count = 0
        try:
            for item, _ in document.iterate_items(page_no=page_no):
                text = getattr(item, "text", None) or getattr(item, "orig", None) or ""
                char_count += len(text)
        except Exception as exc:
            logger.warning(
                "landscape char-count probe failed for page %d (%s)", page_no, exc
            )
            continue
        if char_count < LANDSCAPE_CHAR_THRESHOLD:
            below.append({**p, "char_count": char_count})
    return below


def _rasterize_rotate_page(pdf_path: str, page_no: int, dpi: int = 300) -> str:
    """RFC-035 D2 Phase 2: rasterize a single page at ``dpi`` (fitz already
    applies ``/Rotate`` when rendering) and, if the rendered raster is still
    landscape (width > height — the aspect-ratio case ``/Rotate`` doesn't cover),
    rotate the raster image itself to portrait. Returns the path to a temp PNG.
    Raises on any failure — the caller catches this and falls through to the
    page's original extraction (Design Error Handling item 6)."""
    import io

    import fitz  # PyMuPDF, AGPL-3.0
    from PIL import Image

    with fitz.open(pdf_path) as pdf:
        page = pdf[page_no]
        zoom = dpi / 72
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        png_bytes = pix.tobytes("png")
    img = Image.open(io.BytesIO(png_bytes))
    if img.width > img.height:
        img = img.rotate(-90, expand=True)
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)  # noqa: SIM115
    img.save(tmp.name, format="PNG")
    tmp.close()
    return tmp.name


def _landscape_rasterize_rotate_reextract(
    pdf_path: str, pages: list[dict], ocr_lang_override: list[str] | None = None
) -> list[dict]:
    """RFC-035 D2 Phase 2: for each landscape page flagged below
    ``LANDSCAPE_CHAR_THRESHOLD`` by ``_landscape_pages_below_threshold``,
    rasterize at 300 DPI, rotate to portrait, and re-extract via Docling's
    image pipeline (falling back to local Tesseract OCR if Docling itself
    errors on the rasterized page — HR3-clean, no LLM egress).

    Rasterization/rotation failure logs a warning and skips the page — the
    caller falls through to the page's original (degraded) extraction so
    ``classify_verdict``'s node_count/depth/max_leaf_ratio logic surfaces the
    resulting MARGINAL/FAIL verdict naturally rather than raising (Design
    Error Handling item 6). Routing re-evaluation (feeding recovered
    PictureResults back into flat-mixed classification) is Phase 2's
    follow-on, not this function's concern.
    """
    from .config import ALLOW_AGPL_FALLBACK

    if not ALLOW_AGPL_FALLBACK:
        return []
    results: list[dict] = []
    deadline = time.monotonic() + LANDSCAPE_REEXTRACT_DEADLINE_SECONDS
    for p in pages:
        if len(results) >= MAX_LANDSCAPE_PAGES or time.monotonic() >= deadline:
            logger.warning(
                "landscape reextraction bailing early (%d/%d pages, deadline=%s) for %s",
                len(results),
                MAX_LANDSCAPE_PAGES,
                time.monotonic() >= deadline,
                pdf_path,
            )
            break
        page_no = p["page_no"]
        try:
            png_path = _rasterize_rotate_page(pdf_path, page_no, dpi=300)
        except Exception as exc:
            logger.warning(
                "landscape rasterize/rotate failed for page %d of %s (%s); "
                "falling through to original extraction",
                page_no,
                pdf_path,
                exc,
            )
            continue
        try:
            try:
                converter = _docling_converter(
                    force_full_page_ocr=True,
                    ocr_lang_override=ocr_lang_override,
                    for_image=True,
                )
                result = converter.convert(png_path)
                md = _repair_docling_tables(
                    result.document.export_to_markdown(), doc_name=png_path
                )
                has_pictures = bool(getattr(result.document, "pictures", None))
            except Exception as exc:
                logger.warning(
                    "landscape Docling re-extraction failed for page %d of %s (%s); "
                    "falling back to Tesseract OCR",
                    page_no,
                    pdf_path,
                    exc,
                )
                md = _tesseract_ocr_image(png_path, ocr_lang_override or ["eng"])
                has_pictures = False
        except Exception as exc:
            logger.warning(
                "landscape re-extraction failed for page %d of %s (%s); "
                "falling through to original extraction",
                page_no,
                pdf_path,
                exc,
            )
            continue
        finally:
            with contextlib.suppress(OSError):
                os.unlink(png_path)
        if md and md.strip():
            results.append({"page_no": page_no, "markdown": md, "has_pictures": has_pictures})
    return results


def _recover_picture_text(  # noqa: PLR0915, C901
    pdf_path: str,
    regions: list[dict],
    langs: list[str],
    md: str = "",
    expected_script: str | None = None,
) -> tuple[dict[int, PictureResult], dict[int, str]]:
    """Crop each picture bbox from the PDF, OCR it, and retain the PNG bytes.

    Returns ``{picture_index: PictureResult}`` for every picture region. Each
    result carries ``png_bytes`` (the cropped 300-DPI image), ``ocr_text``
    (Tesseract output, empty if below ``_PICTURE_OCR_MIN_CHARS``), ``page``
    (1-indexed), and ``bbox`` (``{l, t, r, b}``).

    HR3: OCR runs entirely through the LOCAL tesseract binary — no LLM, no
    network egress — so PII rendered inside a chart never leaves the host.

    HR4: this imports ``fitz`` (PyMuPDF, AGPL-3.0) directly for bbox cropping.
    First-party AGPL import on the DEFAULT path; reconciled with the user for
    RFC-015 (2026-07-17). The import is function-scoped and only fires when
    the document actually contains pictures.

    Audit 2026-07-21 findings 10/12: phase 1 crops every valid region SERIALLY
    through one ``fitz.Document`` (PyMuPDF is not shared across threads); phase 2
    OCRs the crops through a bounded ``ThreadPoolExecutor`` (the tesseract CLI is
    a subprocess, safe to parallelize). Decorative gate: when OCR yield is below
    ``_PICTURE_OCR_MIN_CHARS`` the crop's ``png_bytes`` are dropped — unless the
    VLM describe route is enabled downstream, which may still re-mark the image
    as content-bearing via a description.

    D1 (RFC-024): when a region's ``clip_text`` is NOT already contained in
    ``md`` (the Docling markdown export, normalized once here — not per
    region), it is captured directly into ``ocr_text`` (reason
    ``clip_text_captured``) instead of being discarded. This recovers
    chart/infographic text-layer content that Docling misclassified as a
    Picture and that Tesseract OCR on the crop would fail to recognize
    (vector-art labels).

    D5a (RFC-029): skip-gate retention — when the ``page_coverage`` or
    ``clip_text_already_exported`` gate fires, the cropped ``png_bytes`` are
    still captured into the returned ``PictureResult`` so downstream consumers
    (VLM describe route, ``splice_figure_markers``) retain picture context.
    For ``clip_text_already_exported`` the ``clip_text`` is also propagated
    into ``PictureResult.ocr_text`` so ``splice_figure_markers`` can emit a
    ``[Chart text]`` block."""
    from .config import ALLOW_AGPL_FALLBACK, settings

    if not ALLOW_AGPL_FALLBACK:
        logger.warning(
            "picture-region recovery skipped for %s: ALLOW_AGPL_FALLBACK=false "
            "(fitz/PyMuPDF is AGPL-3.0)",
            pdf_path,
        )
        return {}, {}

    import fitz  # PyMuPDF, AGPL-3.0

    md_norm = _normalize_for_containment(md) if _CLIP_TEXT_CAPTURE_ENABLED else ""

    # Phase 1 (serial, single fitz.Document): crop every valid region.
    crops: dict[int, dict] = {}
    clip_captures: dict[int, dict] = {}
    # D5a (RFC-029): regions skipped by a gate but whose png_bytes we still
    # want to retain for downstream context (page_coverage, clip_text_already_exported).
    retained_skips: dict[int, dict] = {}
    skip_reasons: dict[int, str] = {}
    pdf = fitz.open(pdf_path)
    fullpage_ocr_region_count = 0
    try:
        for i, region in enumerate(regions):
            try:
                page_index = region["page"] - 1
                if page_index < 0 or page_index >= pdf.page_count:
                    continue
                page = pdf[page_index]
                rect = _bbox_to_fitz_rect(region["bbox"], page.rect.height, fitz)
                if rect is None:
                    continue
                # D0: skip regions covering >60% of page — full scanned pages, not charts.
                page_area = page.rect.width * page.rect.height
                coverage = (rect.width * rect.height) / page_area if page_area > 0 else 0.0
                if coverage > _PICTURE_PAGE_COVERAGE_THRESHOLD:
                    if _REGION_AWARE_TEXT_CHECK_ENABLED:
                        has_own_text = _region_has_own_text_layer(page, rect)
                        if has_own_text and _TEXT_LAYER_GARBLE_CHECK_ENABLED:
                            region_text = page.get_text("text", clip=rect).strip()
                            if region_text:
                                from .helpers import GarbleContext, check_garble
                                from .script import infer_script

                                if check_garble(
                                    region_text,
                                    expected_script=expected_script or infer_script(region_text),
                                    context=GarbleContext.REGION,
                                ):
                                    has_own_text = False
                    else:
                        has_own_text = _text_layer_has_content(page, expected_script=expected_script)
                    # Reordered: coverage exemption BEFORE MAX_FULLPAGE cap.
                    # A page with no text layer is always exempt (the picture
                    # IS the content) regardless of whether the cap has been
                    # reached — the cap only limits how many exempt pages
                    # actually get the expensive OCR pass.
                    if _COVERAGE_EXEMPT_NO_TEXT_LAYER and not has_own_text:
                        if fullpage_ocr_region_count >= _MAX_FULLPAGE_PICTURE_OCR_REGIONS:
                            logger.warning(
                                "MAX_FULLPAGE_PICTURE_OCR_REGIONS (%d) exceeded for %s; "
                                "skipping further full-page picture exemptions",
                                _MAX_FULLPAGE_PICTURE_OCR_REGIONS,
                                pdf_path,
                            )
                            skip_reasons[i] = "page_coverage"
                            # D5a: retain crop bytes even though OCR is skipped.
                            try:
                                orig_rotation = page.rotation
                                page.set_rotation(0)
                                try:
                                    pix = page.get_pixmap(clip=rect, dpi=300)
                                finally:
                                    page.set_rotation(orig_rotation)
                                retained_skips[i] = {
                                    "png_bytes": pix.tobytes("png"),
                                    "region": region,
                                    "skipped_reason": "page_coverage",
                                }
                            except Exception as _crop_exc:
                                logger.debug(
                                    "D5a: png_bytes crop failed for page_coverage region %d: %s",
                                    i,
                                    _crop_exc,
                                )
                            continue
                        fullpage_ocr_region_count += 1
                        logger.warning(
                            "F1: coverage %.1f%% exceeds threshold but page %d has no text layer; "
                            "exempting from skip (picture IS the page content)",
                            coverage * 100,
                            page_index + 1,
                        )
                    else:
                        skip_reasons[i] = "page_coverage"
                        # D5a: retain crop bytes even though OCR is skipped.
                        try:
                            orig_rotation = page.rotation
                            page.set_rotation(0)
                            try:
                                pix = page.get_pixmap(clip=rect, dpi=300)
                            finally:
                                page.set_rotation(orig_rotation)
                            retained_skips[i] = {
                                "png_bytes": pix.tobytes("png"),
                                "region": region,
                                "skipped_reason": "page_coverage",
                            }
                        except Exception as _crop_exc:
                            logger.debug(
                                "D5a: png_bytes crop failed for page_coverage region %d: %s",
                                i,
                                _crop_exc,
                            )
                        continue
                # D1 (RFC-018/RFC-024): a region with meaningful clip_text is either
                # already exported by Docling (skip) or was misclassified as a
                # Picture and never surfaced (capture), decided by the containment
                # guard against the normalized markdown body.
                clip_text = page.get_text("text", clip=rect).strip()
                if len(clip_text) > _PICTURE_OCR_MIN_CHARS:
                    if _CLIP_TEXT_CAPTURE_ENABLED and not _clip_text_contained(clip_text, md_norm):
                        clip_captures[i] = {
                            "ocr_text": " ".join(clip_text.split()),
                            "region": region,
                        }
                        logger.info(
                            "clip_text_captured for picture region %d in %s (not found in "
                            "Docling markdown export)",
                            i,
                            pdf_path,
                        )
                    else:
                        skip_reason = (
                            "clip_text_already_exported"
                            if _CLIP_TEXT_CAPTURE_ENABLED
                            else "clip_text"
                        )
                        skip_reasons[i] = skip_reason
                        # D5a: for clip_text_already_exported, retain png_bytes AND
                        # propagate clip_text so splice_figure_markers can emit [Chart text].
                        if _CLIP_TEXT_CAPTURE_ENABLED:
                            try:
                                orig_rotation = page.rotation
                                page.set_rotation(0)
                                try:
                                    pix = page.get_pixmap(clip=rect, dpi=300)
                                finally:
                                    page.set_rotation(orig_rotation)
                                retained_skips[i] = {
                                    "png_bytes": pix.tobytes("png"),
                                    "ocr_text": " ".join(clip_text.split()),
                                    "region": region,
                                    "skipped_reason": skip_reason,
                                }
                            except Exception as _crop_exc:
                                logger.debug(
                                    "D5a: png_bytes crop failed for %s region %d: %s",
                                    skip_reason,
                                    i,
                                    _crop_exc,
                                )
                    continue
                # D2 (RFC-023): sub-icon regions (both dims below threshold) are
                # decorative UI glyphs — skip crop+OCR entirely.
                if (
                    rect.width < _DECORATIVE_ICON_MIN_DIM_PT
                    and rect.height < _DECORATIVE_ICON_MIN_DIM_PT
                ):
                    skip_reasons[i] = "decorative_icon"
                    continue
                # D6: zero page rotation before rendering so Tesseract receives a
                # correctly-oriented crop regardless of PDF page-rotation metadata.
                orig_rotation = page.rotation
                page.set_rotation(0)
                try:
                    pix = page.get_pixmap(clip=rect, dpi=300)
                finally:
                    page.set_rotation(orig_rotation)
                crops[i] = {
                    "png_bytes": pix.tobytes("png"),
                    "region": region,
                    "rotation": orig_rotation,
                }
            except Exception as exc:  # D2 (RFC-024): isolate per-region crop failures
                logger.warning(
                    "crop failed for picture region %d in %s (%s); skipping region",
                    i,
                    pdf_path,
                    exc,
                )
                skip_reasons[i] = "crop_error"
                continue
    finally:
        pdf.close()

    recovered: dict[int, PictureResult] = {}
    for i, capture in clip_captures.items():
        bbox = capture["region"]["bbox"]
        recovered[i] = PictureResult(
            ocr_text=capture["ocr_text"],
            page=capture["region"]["page"],
            bbox={"l": bbox.l, "t": bbox.t, "r": bbox.r, "b": bbox.b},
        )
    # D5a (RFC-029): emit retained-skip PictureResults (page_coverage,
    # clip_text_already_exported) so downstream has png_bytes / ocr_text context.
    for i, rs in retained_skips.items():
        bbox = rs["region"]["bbox"]
        pr: PictureResult = PictureResult(
            page=rs["region"]["page"],
            bbox={"l": bbox.l, "t": bbox.t, "r": bbox.r, "b": bbox.b},
            png_bytes=rs["png_bytes"],
            skipped_reason=rs["skipped_reason"],
        )
        if rs.get("ocr_text"):
            pr["ocr_text"] = rs["ocr_text"]
        recovered[i] = pr
    if not crops:
        return recovered, skip_reasons

    def _ocr_one(png_bytes: bytes) -> str:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(png_bytes)
            tmp_path = tmp.name
        try:
            raw = _tesseract_ocr_image(tmp_path, langs)
        finally:
            with contextlib.suppress(OSError):
                os.unlink(tmp_path)
        if len(raw.strip()) > _PICTURE_OCR_MIN_CHARS:
            return " ".join(raw.split())
        return ""

    # Phase 2 (bounded parallel, finding 10): OCR the crops.
    indices = list(crops.keys())
    with ThreadPoolExecutor(max_workers=min(_IMAGE_ENRICH_CONCURRENCY, len(indices))) as pool:
        ocr_texts = dict(
            zip(
                indices,
                pool.map(lambda i: _ocr_one(crops[i]["png_bytes"]), indices),
                strict=True,
            )
        )

    keep_silent_png = settings.vlm_describe_images
    for i in indices:
        region = crops[i]["region"]
        bbox = region["bbox"]
        ocr_text = ocr_texts[i]
        result = PictureResult(
            ocr_text=ocr_text,
            page=region["page"],
            bbox={"l": bbox.l, "t": bbox.t, "r": bbox.r, "b": bbox.b},
        )
        # Finding 12: decorative image (no OCR yield) — drop the crop bytes so
        # no PNG is persisted, unless the VLM route may still describe it.
        if ocr_text or keep_silent_png:
            result["png_bytes"] = crops[i]["png_bytes"]
        # D2 (RFC-023) belt-and-suspenders: a region that passed the size filter
        # but still yielded no OCR text is likely decorative.
        if not ocr_text:
            result["decorative"] = True
        recovered[i] = result
    return recovered, skip_reasons


def _figure_desc_inline(desc: str) -> str:
    """Sanitize a VLM description for the inline ``[Figure: fig-k | desc]`` form
    so it cannot break the single-line ``_FLAT_FIGURE_RE`` grammar."""
    return " ".join(desc.split()).replace("[", "(").replace("]", ")")


def splice_picture_text_for_tree(md: str, pics: list[PictureResult]) -> str:
    """Append OCR text after ``<!-- image -->`` markers for the tree branch.

    Restores the ``_maybe_splice_picture_ocr`` semantics from master that were
    lost when picture OCR was moved to the flat-only path (RFC-020 AD1).

    The ``<!-- image -->`` markers are left **intact** so that the flat branch's
    ``splice_figure_markers`` can still resolve them later if needed.

    Uses ``bind_markers`` from ``picture_plane`` for per-marker alignment
    instead of bailing entirely on a count mismatch.
    """
    from .picture_plane import SkipReason, bind_markers

    if not pics:
        return md
    return bind_markers(md, pics, inject_chart_text=True)


def splice_figure_markers(md: str, pics: list[PictureResult]) -> str:
    """Replace ``<!-- image -->`` markers with ``[Figure: fig-<k>]`` references
    (flat-branch ONLY — audit finding 6: tree-route markdown stays neutral).

    Ordinal matching (RFC-023 D1): the k-th marker is spliced against ``pics[k]``
    when it exists, aligning with ``_enrich_image_blocks``'s ``pic_results[index]``
    lookup (finding 4). Markers count differs from ``len(pics)`` no longer bails
    out — excess markers past ``len(pics)`` (no matching ``PictureResult``) are
    stripped when ``STRIP_SKIPPED_IMAGE_MARKERS=true`` (default), else left as
    neutral markers, matching the existing skipped/decorative behavior below.

    Decorative results (no png/ocr/description — finding 12) keep their neutral
    marker so no unresolvable ``[Figure: fig-k]`` reference is ever emitted.

    Sets ``spliced_into_markdown`` flag on spliced pics instead of the prior
    destructive ``pop('ocr_text')``."""
    from .picture_plane import SkipReason

    if not pics:
        return md
    _LANDSCAPE_REASONS = {SkipReason.LANDSCAPE_FALLBACK.value, "landscape_fallback_picture"}
    real_pics = [p for p in pics if p.get("skipped_reason") not in _LANDSCAPE_REASONS]
    marker_count = md.count(_IMAGE_MARKER)
    if marker_count != len(real_pics):
        logger.warning(
            "figure marker/region count mismatch (%d marker(s) vs %d real picture result(s), "
            "%d landscape fabricated); splicing by ordinal, stripping/neutralizing excess markers",
            marker_count,
            len(real_pics),
            len(pics) - len(real_pics),
        )
    counter = {"i": 0}
    _spliced_indices: set[int] = set()

    def _repl(m: "re.Match[str]") -> str:
        k = counter["i"]
        counter["i"] += 1
        if k >= len(real_pics):
            strip_env = os.environ.get("STRIP_SKIPPED_IMAGE_MARKERS", "true").lower()
            if strip_env != "false":
                return ""
            return m.group(0)
        result = real_pics[k]
        ocr = result.get("ocr_text", "")
        desc = result.get("description", "")
        if not (ocr or desc or result.get("png_bytes")):
            if result.get("skipped_reason") or result.get("decorative"):
                strip_env = os.environ.get("STRIP_SKIPPED_IMAGE_MARKERS", "true").lower()
                if strip_env != "false":
                    return ""
            return m.group(0)
        if desc:
            marker = f"[Figure: fig-{k} | {_figure_desc_inline(desc)}]"
        else:
            marker = f"[Figure: fig-{k}]"
        if ocr:
            _spliced_indices.add(k)
            return marker + "\n\n> [Chart text]: " + ocr
        return marker

    spliced = re.sub(re.escape(_IMAGE_MARKER), _repl, md)
    # Set spliced_into_markdown flag instead of destructive pop('ocr_text')
    for idx in _spliced_indices:
        real_pics[idx]["_spliced_into_markdown"] = True
    return spliced


def _pre_inference_normalize(text: str) -> str:
    """Markdown clean-up run BEFORE heading-depth inference (RFC-015 D5c/D4/D7).

    Ordering is load-bearing: D5c (split run-together headings) must precede D4 (the
    per-line hash-sentinel fix, so ``##Foo ###Bar`` is split before the one-marker-per-
    line pass), which must precede D7 (BiDi reorder) and depth inference (so في is a
    single token by the time the heading regex parses it).

    RFC-029 §1.1: NFKC canonicalization of Arabic Presentation Forms (U+FB50–FDFF,
    U+FE70–FEFF) runs first so all downstream consumers see canonical codepoints.
    The pass is gated on detection — non-Arabic text is untouched (idempotent).
    Design Property 1: NFKC canonicalization idempotence.
    """
    # RFC-029 §1.1 — NFKC only when Arabic Presentation Forms are present.
    # Ranges: Arabic Presentation Forms-A U+FB50–U+FDFF,
    #         Arabic Presentation Forms-B U+FE70–U+FEFF.
    if any("ﭐ" <= ch <= "﷿" or "ﹰ" <= ch <= "﻿" for ch in text):
        text = unicodedata.normalize("NFKC", text)
    text = _split_run_together_headings(text)  # D5c
    text = _fix_fi_hash_substitution(text)  # D4 (moved earlier in the pipeline)
    text = reconstruct_bidi_order(text)  # D7 (Zone-3: sole bidi normalization step)
    return text


def _recover_picture_results(
    md: str,
    document,
    pdf_path: str,
    filename: str | None = None,
    body_for_containment: str | None = None,
    expected_script: str | None = None,
) -> list[PictureResult]:
    """Recover chart/infographic text Docling bucketed into Picture bboxes (RFC-015 D6).

    OCR + crop ONLY — no markdown mutation, no VLM (both moved to the flat branch
    of ``client.index()``, the sole consumer — audit findings 6/8). Gated on
    ``_OCR_ESCALATION`` (mirrors client.py:66) + the presence of a ``<!-- image -->``
    marker, and never fatal.

    Returns a DENSE list: element ``i`` corresponds to the i-th PictureItem in
    ``iterate_items`` order, with an empty ``PictureResult`` placeholder for any
    region whose crop failed — sparse recovery must never shift ordinals
    (finding 4).

    ``body_for_containment``: when provided, the containment check
    (``_normalize_for_containment`` / ``_clip_text_contained``) measures against
    this text instead of ``md``. This fixes the RFC-024 D1 suppression bug where
    ``_document_level_text_fallback`` appends the full pdfium text layer to ``md``
    before containment runs, making every picture's clipped OCR text look "already
    contained" and wrongly skipping legitimate recovery.

    Language detection (RFC-028 D5): ``md`` is the Docling markdown export, which
    is near-empty or all-digits for scanned Arabic PDFs, so ``detect_ocr_langs(md)``
    alone falls through to ``['eng']``. Union with ``detect_ocr_langs(filename)``
    (matching the escalation sites in client.py) so filename script hints survive
    even when the export carries no usable signal."""
    # Zone-6: centralised OCR-mode dispatch replaces ad-hoc boolean gate.
    # Zone-5: per-picture enrichment gate (not page-level garble retry).
    _ocr_mode = decide_ocr_mode(
        ocr_escalation_enabled=_OCR_ESCALATION_PER_PICTURE,
        has_image_markers=_IMAGE_MARKER in md,
    )
    if _ocr_mode == OcrMode.NONE:
        return []
    containment_md = body_for_containment if body_for_containment is not None else md
    try:
        regions = _collect_picture_regions(document)
        if not regions:
            return []
        lang_sources: list[str] = []
        for src in (detect_ocr_langs(filename or ""), detect_ocr_langs(md or "")):
            for lg in src:
                if lg not in lang_sources:
                    lang_sources.append(lg)
        langs = ensure_tessdata(lang_sources)
        recovered, skip_reasons = _recover_picture_text(pdf_path, regions, langs, md=containment_md, expected_script=expected_script)
        if not recovered and not skip_reasons:
            return []
        logger.info(
            "recovered per-picture chart text for %d of %d image(s) in %s",
            len(recovered),
            len(regions),
            pdf_path,
        )
        return [
            recovered.get(i, PictureResult(skipped_reason=skip_reasons.get(i, "unknown")))
            for i in range(len(regions))
        ]
    except Exception as exc:
        logger.warning(
            "per-picture OCR recovery failed for %s (%s); continuing without figures",
            pdf_path,
            exc,
        )
    return []


def _add_vlm_descriptions(pics: list[PictureResult], doc_id: str) -> None:
    """Add VLM-generated descriptions to picture results (HR3-gated, flat-branch only).

    Egress rides ``zdr_egress_gate`` and passes the SAME ``api_base`` the gate
    inspected to ``litellm.completion`` (finding 3). Calls run through a bounded
    ``ThreadPoolExecutor`` (finding 10). Each call is retried once after a short
    backoff; a terminal failure increments ``IMAGE_DESCRIBE_FAILURES`` — matching
    the ``html_to_markdown_with_images._describe`` contract (finding 15)."""
    allowed, api_base = zdr_egress_gate("VLM image descriptions", doc_id=doc_id)
    if not allowed:
        return

    import base64

    from litellm import completion

    from .config import settings
    from .metrics import IMAGE_DESCRIBE_FAILURES

    model = settings.vlm_model
    targets = [(k, pr) for k, pr in enumerate(pics) if pr.get("png_bytes")]
    if not targets:
        return

    def _describe_one(item: tuple[int, PictureResult]) -> None:
        k, result = item
        png_b64 = base64.b64encode(result["png_bytes"]).decode()
        messages = [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{png_b64}",
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "Describe this figure concisely in one sentence. "
                            "Focus on chart type, data series, and key values."
                        ),
                    },
                ],
            }
        ]
        for attempt in (0, 1):
            try:
                resp = completion(
                    model=model,
                    api_base=api_base,
                    messages=messages,
                    max_tokens=150,
                )
                desc = (resp.choices[0].message.content or "").strip()
                if desc:
                    result["description"] = desc
                return
            except Exception as exc:
                if attempt == 0:
                    # Transient failure — retry once after a short backoff.
                    time.sleep(2)
                    continue
                logger.error(
                    "VLM description failed after retry for fig-%d of %s (%s): %s",
                    k,
                    doc_id,
                    type(exc).__name__,
                    str(exc)[:200],
                )
                IMAGE_DESCRIBE_FAILURES.labels(error_type=type(exc).__name__).inc()

    with ThreadPoolExecutor(max_workers=min(_IMAGE_ENRICH_CONCURRENCY, len(targets))) as pool:
        list(pool.map(_describe_one, targets))


# RFC-027 D7: dynamic CHILD_TIMEOUT scaling for the chunked-Docling path. A
# fixed CHILD_TIMEOUT sized for a single-pass conversion is what oversized
# PDFs die to in the first place; the chunked path needs a timeout budget
# proportional to how many independent Docling passes it runs.
_CHUNKED_DOCLING_BASE_TIMEOUT_S = 300
# RFC-028 D0: 600 -> 1500. The prior constant made chunked_docling_timeout_s(2)
# (1500s) *lower* than the fixed CHILD_TIMEOUT (1770s) it was meant to extend,
# so wiring it in without raising this would have shrunk the timeout budget for
# world-stats-pocketbook-2023.pdf (292 pages, observed 24-49min conversion).
_CHUNKED_DOCLING_PER_CHUNK_TIMEOUT_S = 1500


def chunked_docling_timeout_s(chunk_count: int) -> int:
    """RFC-027 D7: ``base_timeout + (chunk_count * per_chunk_timeout)``.

    Consumed by the worker's per-job CHILD_TIMEOUT so a chunked conversion
    gets a budget proportional to how many independent Docling passes it
    runs, instead of the fixed single-pass timeout that oversized PDFs die to.
    """
    return _CHUNKED_DOCLING_BASE_TIMEOUT_S + chunk_count * _CHUNKED_DOCLING_PER_CHUNK_TIMEOUT_S


try:
    from pdf_inspector import detect_pdf as _detect_pdf

    _pdf_inspector_available = True
except ImportError:
    _pdf_inspector_available = False
    _detect_pdf = None  # type: ignore[assignment]


def _run_pdf_inspector(pdf_path: str) -> dict | None:
    """Run pdf-inspector classification and return a dict, or None on failure."""
    if not _pdf_inspector_available:
        return None
    try:
        t0 = time.monotonic()
        result = _detect_pdf(pdf_path)
        elapsed = time.monotonic() - t0
        from .metrics import PDF_INSPECTOR_CLASSIFICATIONS, PDF_INSPECTOR_LATENCY

        PDF_INSPECTOR_LATENCY.observe(elapsed)
        PDF_INSPECTOR_CLASSIFICATIONS.labels(pdf_type=result.pdf_type).inc()
        return {
            "pdf_type": result.pdf_type,
            "confidence": result.confidence,
            "pages_needing_ocr": list(result.pages_needing_ocr),
            "has_encoding_issues": getattr(result, "has_encoding_issues", False),
        }
    except Exception:
        logging.getLogger(__name__).debug(
            "pdf-inspector classify failed for %s", pdf_path, exc_info=True
        )
        return None


def probe_conversion_route(pdf_path: str) -> tuple[int, bool, dict | None]:
    """RFC-028 D0: cheap pre-flight probe run by ``converters_cli`` before the
    heavy conversion pipeline starts, so the worker can size its child timeout
    from the child's own startup handshake instead of re-deriving page count
    independently (which risks worker/child disagreement on a page-count failure).

    Returns ``(chunk_count, is_docling_route, pdf_classification)`` using the
    same pymupdf page-count read and ``MAX_DOCLING_PAGES`` threshold as the
    routing guard at the top of ``pdf_to_markdown_docling``.  Non-PDF inputs
    and PDFs whose page count cannot be read report ``is_docling_route=False``
    so the worker falls back to the fixed ``CHILD_TIMEOUT`` unconditionally.

    ``pdf_classification`` is a dict with pdf-inspector shadow-mode results
    (pdf_type, confidence, pages_needing_ocr, has_encoding_issues), or None
    when pdf-inspector is not installed or classification fails.  Shadow mode:
    classification is logged and metered. When PDF_INSPECTOR_PRECLASSIFY=1
    (config.py), the classification influences behavior: scanned/image-based
    documents with confidence >= 0.90 force first-pass OCR (client.py) and
    receive a 16.5x timeout multiplier (worker.py). When the flag is disabled
    (default), classification is shadow-mode only.
    """
    if not pdf_path.lower().endswith(".pdf"):
        return 1, False, None
    from .config import MAX_DOCLING_PAGES

    classification = _run_pdf_inspector(pdf_path)

    try:
        import pypdfium2 as pdfium  # BSD-3/Apache-2, not fitz/PyMuPDF (AGPL-3.0)

        pdoc = pdfium.PdfDocument(pdf_path)
        try:
            page_count = len(pdoc)
        finally:
            pdoc.close()
    except Exception:
        return 1, False, classification
    if page_count <= 0:
        return 1, False, classification
    if MAX_DOCLING_PAGES > 0 and page_count > MAX_DOCLING_PAGES:
        return math.ceil(page_count / MAX_DOCLING_PAGES), True, classification
    return 1, True, classification


def _repair_docling_tables(md: str, doc_name: str = "") -> str:
    """RFC-029 D4 (Task 5.1, Property 6) — post-export table-repair pass.

    Runs after every Docling ``export_to_markdown()`` call to correct two
    systematic Docling GFM rendering artefacts:

    1. **Degenerate duplicate-cell rows**: pipe-table data rows where EVERY
       non-separator cell is byte-identical (e.g. Docling emitting the same
       cell value repeated across all columns due to table-cell merging
       ambiguity).  A row is collapsed to a single ``| value |`` cell only
       when the identical-cell count exceeds ``_RFC029_TABLE_MIN_COLLAPSE_COLS``
       (default 3) — avoids collapsing legitimate 1- or 2-col tables that
       happen to share a value across columns.

    2. **GFM-aligned whitespace padding**: Docling right-pads every pipe-table
       cell to column-width for visual alignment.  The downstream tree builder
       and flat-table parser both strip whitespace, so the padding is harmless
       for semantics but inflates character counts (up to ~10x for wide
       statistical tables).  Re-emitting with single-space padding recovers the
       inflation without data loss.

    Both transforms are heuristic-only, work on the raw markdown string, and
    require no external dependencies (stdlib ``re`` only).  When
    ``_RFC029_TABLE_DEDUP_ENABLED`` is falsy the function is a no-op.

    Content-preservation: collapsed rows replace the original row text with a
    single-cell row; non-collapsed rows are re-emitted with stripped (single-
    space-padded) cell content, preserving every non-whitespace character.
    Separator rows (``|---|``) are re-emitted as ``| --- |`` (minimal form).

    RFC-034 D10 Phase A: logs before/after char counts plus collapsed-row and
    whitespace-stripped-char counts (read-only diagnostic for Phase B).
    """
    if not _RFC029_TABLE_DEDUP_ENABLED or not md:
        return md

    chars_before = len(md)
    collapsed_rows = 0
    whitespace_stripped = 0
    lines = md.split("\n")
    out: list[str] = []
    prev_was_separator = False

    for line in lines:
        stripped = line.strip()
        # Only process lines that look like pipe-table rows.
        if not stripped.startswith("|") or not stripped.endswith("|"):
            out.append(line)
            continue

        # Split on pipe, drop leading/trailing empty strings from the outer | |.
        raw_cells = stripped.split("|")
        cells = [c.strip() for c in raw_cells[1:-1]]

        if not cells:
            out.append(line)
            continue

        # Detect separator row (cells contain only dashes, colons, spaces).
        if all(
            set(c.replace("-", "").replace(":", "").replace(" ", "")) == set() and c for c in cells
        ):
            # Re-emit in minimal form: | --- | --- | ...
            out.append("| " + " | ".join("---" for _ in cells) + " |")
            prev_was_separator = True
            continue

        # Check for all-identical degenerate row.
        unique_vals = set(cells)
        if len(unique_vals) == 1 and len(cells) > _RFC029_TABLE_MIN_COLLAPSE_COLS:
            # RFC-034 D17: mixed-script rows (Arabic + Latin) are likely
            # legitimate bilingual data, not a Docling merge artefact --
            # skip the collapse and re-emit the row unchanged.
            cell_text = cells[0]
            has_arabic = any(_is_arabic_char(c) for c in cell_text)
            has_latin = bool(re.search(r"[A-Za-z]", cell_text))
            if has_arabic and has_latin:
                new_line = "| " + " | ".join(cells) + " |"
                whitespace_stripped += max(0, len(line) - len(new_line))
                out.append(new_line)
                prev_was_separator = False
                continue
            # RFC-035 D0: the row immediately after a separator is the first
            # body row, not a Docling merge artefact -- repeated labels here
            # (e.g. a sub-header row) are structural. Skip the collapse.
            if prev_was_separator:
                new_line = "| " + " | ".join(cells) + " |"
                whitespace_stripped += max(0, len(line) - len(new_line))
                out.append(new_line)
                prev_was_separator = False
                continue
            # Collapse: emit a single cell with the shared value.
            collapsed_rows += 1
            out.append("| " + cells[0] + " |")
            prev_was_separator = False
            continue

        # Normal row: re-emit with minimal single-space padding (strips GFM alignment).
        new_line = "| " + " | ".join(cells) + " |"
        whitespace_stripped += max(0, len(line) - len(new_line))
        out.append(new_line)
        prev_was_separator = False

    result = "\n".join(out)
    logger.info(
        "table_repair: %s chars %d->%d, collapsed_rows=%d, whitespace_stripped=%d",
        doc_name,
        chars_before,
        len(result),
        collapsed_rows,
        whitespace_stripped,
    )
    return result


def _docling_chunk_worker(
    result_queue: "multiprocessing.Queue",
    pdf_path: str,
    force_full_page_ocr: bool,
    ocr_lang_override: list[str] | None,
    expected_script: str | None = None,
) -> None:
    """Run ``pdf_to_markdown_docling`` in a child process (D0 fix).

    Executed as the target of a ``multiprocessing.Process`` so the parent can
    ``terminate()`` it on timeout and guarantee the work actually stops --
    unlike a ``ThreadPoolExecutor`` thread, which keeps running past
    ``future.result(timeout=...)`` because that only abandons the wait.
    """
    try:
        result_queue.put(("ok", pdf_to_markdown_docling(
            pdf_path,
            force_full_page_ocr=force_full_page_ocr,
            ocr_lang_override=ocr_lang_override,
            expected_script=expected_script,
        )))
    except Exception as exc:  # noqa: BLE001 -- re-raised in parent
        try:
            result_queue.put(("error", exc))
        except Exception:  # exc itself unpicklable -- send a picklable stand-in
            result_queue.put(
                ("error", RuntimeError(f"{type(exc).__name__}: {exc}"))
            )


def _run_docling_chunk_with_timeout(
    pdf_path: str,
    *,
    force_full_page_ocr: bool,
    ocr_lang_override: list[str] | None,
    timeout_s: float,
    expected_script: str | None = None,
) -> tuple[str, list[PictureResult]]:
    """Run one Docling chunk conversion in a killable subprocess (D0 fix).

    Replaces the plain ``ThreadPoolExecutor`` used previously: a
    ``multiprocessing.Process`` can be ``terminate()``-d on timeout, which
    actually stops the in-flight Docling work rather than merely abandoning
    the wait for it. This lets the arq worker's child process exit cleanly
    within its own timeout budget instead of surviving to ``JOB_TIMEOUT``.
    """
    ctx = multiprocessing.get_context("spawn")
    result_queue: multiprocessing.Queue = ctx.Queue()
    proc = ctx.Process(
        target=_docling_chunk_worker,
        args=(result_queue, pdf_path, force_full_page_ocr, ocr_lang_override, expected_script),
        daemon=True,
    )
    proc.start()
    # Drain the queue BEFORE join()ing: a large result (markdown +
    # PictureResult png_bytes) exceeds the queue's pipe buffer, and the child
    # cannot exit until the parent reads it -- join-first would deadlock until
    # the timeout and misreport a *successful* chunk as timed out. The poll
    # loop also detects a child that died without reporting (native segfault
    # in Docling/OCR), which a bare blocking get() would hang on forever.
    deadline = time.monotonic() + timeout_s
    outcome: tuple[str, object] | None = None
    while outcome is None:
        try:
            outcome = result_queue.get(timeout=1.0)
        except queue_mod.Empty:
            if time.monotonic() >= deadline:
                break
            if not proc.is_alive():
                # Child exited; give the queue feeder one final grace read in
                # case the result landed between the Empty and the liveness
                # check, then treat silence as a crash.
                try:
                    outcome = result_queue.get(timeout=1.0)
                except queue_mod.Empty:
                    break
    if outcome is None:
        if proc.is_alive():
            proc.terminate()
            proc.join(5)
            if proc.is_alive():
                proc.kill()
                proc.join()
            raise FuturesTimeoutError(
                f"Docling chunk timed out after {timeout_s}s: {pdf_path}"
            )
        raise RuntimeError(
            "Docling chunk worker died without a result "
            f"(exitcode={proc.exitcode}): {pdf_path}"
        )
    proc.join(5)
    if proc.is_alive():  # lingering after reporting -- reap it
        proc.kill()
        proc.join()
    status, payload = outcome
    if status == "error":
        raise cast(Exception, payload)
    return cast("tuple[str, list[PictureResult], dict[str, dict]]", payload)


def _pdf_to_markdown_docling_chunked(
    pdf_path: str,
    page_count: int,
    max_pages: int,
    force_full_page_ocr: bool = False,
    ocr_lang_override: list[str] | None = None,
    expected_script: str | None = None,
) -> tuple[str, list[PictureResult], dict[str, dict]]:
    """RFC-027 D7 chunked-Docling route for PDFs exceeding MAX_DOCLING_PAGES.

    Splits ``pdf_path`` into ``ceil(page_count / max_pages)`` page-boundary
    chunks via ``pymupdf`` (``fitz``) -- the project's single PDF-primitive
    layer, no ``pymupdf4llm`` (CLAUDE.md Hard Rule 4) -- runs each chunk
    through the existing standard ``pdf_to_markdown_docling`` pipeline
    independently, and concatenates the resulting markdown. Each chunk's page
    count is <= ``max_pages`` by construction, so the recursive call takes the
    direct single-pass route rather than re-entering this function.

    Minor heading-level discontinuities at chunk joins are an accepted
    trade-off (RFC-027 D7 risk acceptance) -- the downstream tree-building
    ``_relevel_by_containment`` pass normalizes heading depth across the
    concatenated output.
    """
    from .config import ALLOW_AGPL_FALLBACK

    if not ALLOW_AGPL_FALLBACK:
        raise RuntimeError(
            f"cannot chunk {pdf_path} for the oversized-PDF route: fitz "
            "(PyMuPDF, AGPL-3.0) is required and ALLOW_AGPL_FALLBACK=false"
        )
    import fitz  # PyMuPDF

    chunk_count = math.ceil(page_count / max_pages)
    logger.info(
        "chunked-Docling route: %s (%d pages) -> %d chunk(s) of <= %d pages",
        pdf_path,
        page_count,
        chunk_count,
        max_pages,
    )
    src = fitz.open(pdf_path)
    md_parts: list[str] = []
    pic_results: list[PictureResult] = []
    try:
        for i in range(chunk_count):
            start = i * max_pages
            end = min(start + max_pages, page_count)
            # SIM115 rationale: the temp FILE must outlive this statement -- it is
            # written, then re-opened by name below and unlinked in `finally`. A
            # context manager would close/delete it before it is ever used.
            tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)  # noqa: SIM115
            tmp.close()
            try:
                writer = fitz.open()
                try:
                    writer.insert_pdf(src, from_page=start, to_page=end - 1)
                    writer.save(tmp.name)
                finally:
                    writer.close()
                try:
                    chunk_md, chunk_pics, _chunk_stages = _run_docling_chunk_with_timeout(
                        tmp.name,
                        force_full_page_ocr=force_full_page_ocr,
                        ocr_lang_override=ocr_lang_override,
                        timeout_s=_CHUNKED_DOCLING_PER_CHUNK_TIMEOUT_S,
                        expected_script=expected_script,
                    )
                except FuturesTimeoutError:
                    # RFC-027 D7: an individually heavy chunk still times out on the
                    # Docling pipeline -- fall back to pymupdf text-layer-only
                    # extraction (no tables/figures) rather than losing the chunk
                    # entirely. No pymupdf4llm (CLAUDE.md Hard Rule 4). The document
                    # lands MARGINAL downstream due to the resulting flat structure.
                    logger.warning(
                        "chunk %d/%d of %s timed out on Docling; falling back to "
                        "pymupdf text-layer extraction",
                        i + 1,
                        chunk_count,
                        pdf_path,
                    )
                    chunk_doc = fitz.open(tmp.name)
                    try:
                        chunk_md = "\n\n".join(page.get_text() or "" for page in chunk_doc)
                    finally:
                        chunk_doc.close()
                    chunk_pics = []
            finally:
                with contextlib.suppress(OSError):
                    os.unlink(tmp.name)
            md_parts.append(chunk_md)
            for pic in chunk_pics:
                # Re-base chunk-relative page numbers to document-level pages so
                # the persisted PictureResult metadata (client.py block["page"])
                # stays correct for chunks after the first.
                if "page" in pic:
                    pic["page"] = pic["page"] + start
            pic_results.extend(chunk_pics)
    finally:
        src.close()
    # Per-chunk stage tables are not merged -- out of scope for Zone 4 initial
    # landing. extraction_stages is empty for chunked/oversized PDFs.
    return "\n\n".join(md_parts), pic_results, {}


def _build_candidate(md: str) -> str:
    """Normalise a candidate markdown source BEFORE heading-depth inference.

    Ordering matters: Arabic structural headings must be injected before
    _pre_inference_normalize runs its NFKC fold + bidi reconstruction, because
    the injection regex matches raw Arabic text that NFKC would alter. German
    clause and English article headings follow, then the pipeline-level
    normalize pass that splits run-together headings, fixes fi-hash
    substitutions, and reconstructs bidi order.
    """
    md = _inject_arabic_structural_headings(md)
    md = _inject_german_clause_headings(md)
    md = _inject_english_article_headings(md)
    md = _pre_inference_normalize(md)
    return md


@dataclasses.dataclass
class StageRecord:
    """Per-stage provenance entry for the extraction pipeline."""

    name: str
    chars_before: int
    chars_after: int
    char_delta: int
    headings_before: int
    headings_after: int
    heading_delta: int
    error: str | None = None


@dataclasses.dataclass(frozen=True)
class Candidate:
    """Immutable (md, heading_pages) pair — keeps the two values that describe
    a single pipeline candidate in lock-step so they can never drift apart."""

    md: str
    heading_pages: dict[str, list[int]] = dataclasses.field(default_factory=dict)


def _run_stages(
    md: str, stages: list[tuple[str, Callable[[str], str]]]
) -> tuple[str, dict[str, dict]]:
    """Run a sequence of string-mutating stages with per-stage provenance.

    Each ``(name, fn)`` pair is called independently: a failure in stage N
    does not skip stages N+1..last. On failure the stage's error is recorded
    and ``md`` is left unchanged for that stage.

    Returns ``(md, records)`` where ``records`` is a name-keyed dict — stage
    names are unique per call.
    """
    records: dict[str, dict] = {}
    for name, fn in stages:
        chars_before = len(md)
        headings_before = len(_HEADING_RE.findall(md))
        try:
            result = fn(md)
            chars_after = len(result)
            headings_after = len(_HEADING_RE.findall(result))
            records[name] = dataclasses.asdict(
                StageRecord(
                    name=name,
                    chars_before=chars_before,
                    chars_after=chars_after,
                    char_delta=chars_after - chars_before,
                    headings_before=headings_before,
                    headings_after=headings_after,
                    heading_delta=headings_after - headings_before,
                )
            )
            md = result
        except Exception as exc:
            logger.warning("extraction stage %r failed: %s", name, exc)
            records[name] = dataclasses.asdict(
                StageRecord(
                    name=name,
                    chars_before=chars_before,
                    chars_after=chars_before,
                    char_delta=0,
                    headings_before=headings_before,
                    headings_after=headings_before,
                    heading_delta=0,
                    error=str(exc),
                )
            )
    return md, records


def pdf_to_markdown_docling(  # noqa: PLR0915, C901
    pdf_path: str,
    force_full_page_ocr: bool = False,
    ocr_lang_override: list[str] | None = None,
    max_pages: int | None = None,
    expected_script: str | None = None,
) -> tuple[str, list[PictureResult], dict[str, dict]]:
    """MIT-licensed layout-aware PDF route (RFC-003 D3 / HR4 AGPL escape).

    Returns ``(markdown, pic_results, extraction_stages)``. The markdown keeps
    bare ``<!-- image -->`` markers (no figure references — audit finding 6);
    ``pic_results[i]`` corresponds
    to the i-th PictureItem in ``iterate_items`` order and always has
    ``len == number of picture regions`` when non-empty (dense — finding 4).

    Docling's Heron RT-DETRv2 layout model + TableFormer -> markdown -> relevel
    headings -> normalize dashes. Validated head-to-head against pymupdf4llm on
    the German insurance corpus (2026-05-31): Docling resolves the ``fl``-ligature
    corruption pymupdf4llm leaves in legal terms (e.g. ``Haftpflicht`` rendered as
    ``Haftpficht``), at ~2.5-6x the CPU runtime.

    The accelerator is pinned to CPU unconditionally — no MPS, no CUDA. This is a
    deliberate operational choice (everything runs on CPU for now) and also sidesteps
    the Apple-MPS crash: transformers' ``rt_detr_v2`` hardcodes float64 in its sin/cos
    position embedding, which MPS rejects (the same wall poc-insurance-chat's
    ``_resolve_accelerator_device`` works around by coercing to CPU on darwin).

    OCR, when enabled, runs through the installed Tesseract binary (CLI engine) so
    the system ``deu``/``eng`` language data is used; point ``TESSDATA_PREFIX`` at the
    directory holding ``deu.traineddata`` (e.g. the repo-local ``.tessdata/``).
    Env knobs:
      ``DOCLING_DO_OCR``   1|0 (default 0 — text-layer PDFs need no OCR)
      ``DOCLING_OCR_LANG`` comma list (default ``deu,eng``) when OCR is on
      ``DOCLING_ARTIFACTS_PATH`` dir of pre-downloaded model weights for offline use
        (set in the container image; unset locally -> weights fetched from HF on first use)

    Raises on empty extraction so the caller falls back to the next converter.
    """
    # RFC-027 D7: oversized PDFs die to CHILD_TIMEOUT on a single direct-conversion
    # pass. Guard the page count via pymupdf (no pymupdf4llm -- CLAUDE.md Hard
    # Rule 4) before touching the Docling converter and route to the chunked
    # path instead.
    from .config import ALLOW_AGPL_FALLBACK, MAX_DOCLING_PAGES

    effective_max_pages = max_pages if max_pages is not None else MAX_DOCLING_PAGES
    if not ALLOW_AGPL_FALLBACK:
        logger.warning(
            "chunked-Docling page-count guard skipped for %s: ALLOW_AGPL_FALLBACK=false "
            "(fitz/PyMuPDF is AGPL-3.0)",
            pdf_path,
        )
        page_count = 0
    else:
        try:
            import fitz  # PyMuPDF

            with fitz.open(pdf_path) as doc:
                page_count = doc.page_count
        except Exception as exc:
            logger.warning(
                "could not read page count for %s (%s); skipping chunked-Docling guard",
                pdf_path,
                exc,
            )
            page_count = 0
    if effective_max_pages > 0 and page_count > effective_max_pages:
        return _pdf_to_markdown_docling_chunked(
            pdf_path,
            page_count=page_count,
            max_pages=effective_max_pages,
            force_full_page_ocr=force_full_page_ocr,
            ocr_lang_override=ocr_lang_override,
            expected_script=expected_script,
        )

    # Reuse the process-cached converter (see _docling_converter): a fresh
    # DocumentConverter per call leaks ~250 MB/doc that torch never frees.
    converter = _docling_converter(
        force_full_page_ocr=force_full_page_ocr, ocr_lang_override=ocr_lang_override
    )
    # RFC-035 D2 Phase 1: read-only landscape probe, tags pages for the future
    # rasterize-rotate-reextract fallback (Phase 2). Does not alter extraction.
    landscape_pages = _tag_landscape_pages_for_fallback(pdf_path)
    if any(p["is_landscape"] for p in landscape_pages):
        logger.info(
            "landscape pages detected in %s: %s",
            pdf_path,
            [p["page_no"] for p in landscape_pages if p["is_landscape"]],
        )

    # RFC-026 D2: normalize per-page rotation before extraction so landscape/
    # rotated pages get correct coordinate mapping instead of fragmenting text
    # into near-empty nodes. Returns pdf_path unchanged when no page needs it.
    docling_input_path = _normalize_pdf_page_rotation(pdf_path)
    try:
        result = converter.convert(docling_input_path)
    finally:
        if docling_input_path != pdf_path:
            with contextlib.suppress(OSError):
                os.unlink(docling_input_path)

    # RFC-035 D2 Phase 2 trigger: for pages tagged landscape above, compare the
    # primary extraction's char count against LANDSCAPE_CHAR_THRESHOLD. Detection
    # only here — the rasterize-rotate-reextract fallback itself is Phase 2 proper.
    landscape_below_threshold = _landscape_pages_below_threshold(
        result.document, landscape_pages
    )
    landscape_fallback_pages: list[dict] = []
    if landscape_below_threshold:
        logger.info(
            "landscape pages below LANDSCAPE_CHAR_THRESHOLD (%d chars) in %s: %s",
            LANDSCAPE_CHAR_THRESHOLD,
            pdf_path,
            [(p["page_no"], p["char_count"]) for p in landscape_below_threshold],
        )
        # RFC-035 D2 Phase 2: rasterize-rotate-reextract fallback. Never fatal —
        # a failure here falls through to the original (degraded) extraction and
        # classify_verdict surfaces the resulting MARGINAL/FAIL verdict naturally.
        landscape_fallback_pages = _landscape_rasterize_rotate_reextract(
            pdf_path, landscape_below_threshold, ocr_lang_override=ocr_lang_override
        )
        if landscape_fallback_pages:
            logger.info(
                "landscape rasterize-rotate-reextract recovered %d page(s) for %s: %s",
                len(landscape_fallback_pages),
                pdf_path,
                [p["page_no"] for p in landscape_fallback_pages],
            )

    # Capture the RAW Docling markdown BEFORE the add-on runs: ResultPostprocessor
    # mutates result.document in place (it demotes unmatched headings to body text),
    # so this is the only chance to retain the full heading set for the Rank-1
    # over-prune fallback below.
    raw_md = _repair_docling_tables(result.document.export_to_markdown(), doc_name=pdf_path)

    # Snapshot heading -> [page_no, ...] from the RAW (pre-add-on) document: the
    # add-on demotes unmatched headings to body text in place, so this is the only
    # chance to retain page provenance for the over-prune raw_md fallback path. The
    # outline depth-recovery step below (used only for numberless flat-prose docs)
    # maps rendered headings to PDF-outline sections BY this page.
    try:
        heading_pages_raw = _collect_heading_pages(result.document)
    except Exception as exc:
        logger.warning("could not collect raw heading pages for %s (%s)", pdf_path, exc)
        heading_pages_raw = {}

    # docling-hierarchical-pdf (krrome) rebuilds heading SELECTION from the PDF
    # outline/numbering, dropping the font-size false positives Docling otherwise
    # emits as headings (page numbers, letter-spaced body text, clause fragments).
    # Validated on the German corpus 2026-05-31: cuts noisy headings 34-94%.
    # Optional + third-party (single-maintainer) — never let it break ingestion;
    # degrade to raw Docling headings on any failure.
    try:
        from hierarchical.postprocessor import ResultPostprocessor

        # Rank-2: teach the add-on to tolerate publisher numbering prefixes (the
        # TOC title omits the in-document "BHB N"/"A."/"I." prefix) BEFORE it runs,
        # so it keeps the real headings instead of demoting them. Guarded + never
        # fatal — the Rank-1 fallback below covers any patch failure.
        try:
            _patch_hierarchical_infer()
        except Exception as exc:
            logger.warning(
                "could not patch hierarchical infer() (%s); relying on raw-docling fallback",
                exc,
            )
        ResultPostprocessor(result, source=pdf_path).process()
    except ImportError:
        logger.warning(
            "docling-hierarchical-pdf not installed; using raw docling headings. "
            "Install it to recover clean heading selection."
        )
    except Exception as exc:
        logger.warning(
            "hierarchical add-on postprocess failed for %s (%s); using raw docling headings",
            pdf_path,
            exc,
        )

    # Re-promote the deep numbered clauses the add-on demoted to body text
    # (e.g. AKB "A.1.1"/"A.1.1.1"), restoring the tree depth the add-on prunes.
    # Same defensive contract as the add-on: re-promotion must NEVER be fatal —
    # on any failure degrade to the add-on's selection.
    try:
        n_promo = _repromote_numbered_headings(result.document)
        if n_promo > 0:
            logger.info(
                "re-promoted %d demoted numbered clause(s) to headings for %s",
                n_promo,
                pdf_path,
            )
    except Exception as exc:
        logger.warning(
            "heading re-promotion failed for %s (%s); using add-on selection",
            pdf_path,
            exc,
        )

    post_md = _repair_docling_tables(result.document.export_to_markdown(), doc_name=pdf_path)
    if not post_md or not post_md.strip():
        raise RuntimeError(f"docling produced empty output for {pdf_path}")

    extraction_stages: dict[str, dict] = {}

    # Provenance: docling convert (non-string-mutation, manual entry).
    extraction_stages["docling_convert"] = {
        "name": "docling_convert",
        "chars_before": 0,
        "chars_after": len(raw_md),
        "char_delta": len(raw_md),
        "headings_before": 0,
        "headings_after": len(_HEADING_RE.findall(raw_md)),
        "heading_delta": len(_HEADING_RE.findall(raw_md)),
        "error": None,
    }

    # Provenance: hierarchical add-on (non-string-mutation, manual entry).
    extraction_stages["hierarchical_addon"] = {
        "name": "hierarchical_addon",
        "chars_before": len(raw_md),
        "chars_after": len(post_md),
        "char_delta": len(post_md) - len(raw_md),
        "headings_before": len(_HEADING_RE.findall(raw_md)),
        "headings_after": len(_HEADING_RE.findall(post_md)),
        "heading_delta": len(_HEADING_RE.findall(post_md)) - len(_HEADING_RE.findall(raw_md)),
        "error": None,
    }

    # Page map for the post-add-on candidate's outline step (the RAW map captured
    # before the add-on is used for the raw candidate, keeping each map in sync with
    # the markdown it relevels — see _collect_heading_pages).
    try:
        heading_pages_post = _collect_heading_pages(result.document)
    except Exception as exc:
        logger.warning("could not collect post-add-on heading pages for %s (%s)", pdf_path, exc)
        heading_pages_post = {}

    # Build immutable Candidate pairs (md + heading_pages) via the unified
    # _candidate_from_document entry point so the two values never drift.
    post_candidate = _candidate_from_document(post_md, heading_pages_post, pdf_path)
    raw_candidate = _candidate_from_document(raw_md, heading_pages_raw, pdf_path)

    post_headings = len(_HEADING_RE.findall(post_candidate.md))
    raw_headings = len(_HEADING_RE.findall(raw_candidate.md))

    # Gate-aware source selection (HR5 / over-prune). Recover depth on the CLEANER
    # post-add-on markdown first; if that tree would still fail the structural gate
    # (node_count<3 or depth<2) but the RICHER raw Docling markdown recovers a valid
    # tree, use raw. This subsumes the old `post<3<=raw` count guard AND catches
    # PROPORTIONAL pruning the count guard missed: e.g. Hundehalter/Pferdehalter-
    # haftpflicht, where the add-on demotes ~128 numbered headings to 4 flat ones —
    # 4 is not <3 so the count guard never fired, yet raw_md's numbering chain
    # recovers real depth. raw Docling is ligature-correct + MIT (HR4). The real
    # gate (validate_tree) still runs downstream; this only picks the better source.
    selected = post_candidate
    if not _has_structural_depth(post_candidate.md) and raw_headings >= 3 and raw_headings > post_headings:
        if _has_structural_depth(raw_candidate.md):
            logger.warning(
                "post-add-on tree failed the structural gate (%d heading(s), max-level %d) "
                "for %s; using raw docling markdown (%d headings)",
                post_headings,
                _max_heading_level(post_candidate.md),
                pdf_path,
                raw_headings,
            )
            selected = raw_candidate

    # Runtime contract: the selected source must be a Candidate so the
    # downstream pipeline can rely on .md / .heading_pages being present
    # and the pair being frozen (immutable).
    if not isinstance(selected, Candidate):
        raise TypeError(
            f"source selection must yield a Candidate, got {type(selected).__name__}"
        )

    md = selected.md
    heading_pages_for_md = selected.heading_pages

    # String-mutating stages, each independently fail-open with provenance.
    # Split into pre-fallback and post-fallback so body_for_containment is
    # captured between them (RFC-024 D1 suppression fix).
    pre_fallback_stages: list[tuple[str, Callable[[str], str]]] = [
        ("normalize_indented_headings", _normalize_indented_headings),
    ]
    md, pre_records = _run_stages(md, pre_fallback_stages)
    extraction_stages.update(pre_records)

    # RFC-024 D1 fix: snapshot md BEFORE _document_level_text_fallback appends
    # the raw pdfium text layer, so the containment check in picture recovery
    # measures against genuine Docling-exported content only.
    body_for_containment = md

    post_fallback_stages: list[tuple[str, Callable[[str], str]]] = [
        ("document_level_text_fallback", functools.partial(_document_level_text_fallback, pdf_path=pdf_path, expected_script=expected_script)),
        ("splice_landscape_fallback", functools.partial(
            _splice_landscape_fallback,
            landscape_fallback_pages=landscape_fallback_pages,
            heading_pages=heading_pages_for_md,
        )),
    ]
    md, post_records = _run_stages(md, post_fallback_stages)
    extraction_stages.update(post_records)

    # Audit findings 1/6/11: picture results travel UP THE CALL STACK as part of
    # the return value (a thread-local set on the to_thread pool thread was
    # invisible to the event loop and pinned crop bytes for the process life).
    # The markdown keeps neutral `<!-- image -->` markers — the [Figure: fig-N]
    # splice and the VLM describe step run only in client.index()'s flat branch.
    pic_results = _recover_picture_results(
        md, result.document, pdf_path, os.path.basename(pdf_path),
        body_for_containment=body_for_containment,
        expected_script=expected_script,
    )
    # RFC-035 D2 Fix (Routing interaction / task-5-4): the rasterize-rotate-
    # reextract fallback re-runs Docling on a standalone rasterized page image,
    # so any Picture regions it detects live in that re-extraction's own
    # result.document, not the primary result.document _recover_picture_results
    # just scanned. Surface a routing-only marker per fallback page that fired
    # Docling picture detection so client.index() can re-evaluate classification/
    # routing the same way the portrait companion does (Design Property 3) —
    # this is a signal for the flat-mixed routing decision, not a content-
    # bearing crop, so it deliberately carries no ocr_text/png_bytes and is
    # inert to splice_figure_markers' marker-count alignment (degrades to
    # neutral on mismatch per that function's docstring).
    for p in landscape_fallback_pages:
        if p.get("has_pictures"):
            pic_results.append(
                PictureResult(page=p["page_no"], skipped_reason=SkipReason.LANDSCAPE_FALLBACK.value)
            )

    # Provenance: picture recovery (non-string-mutation, manual entry).
    recovered_count = sum(1 for pr in pic_results if pr.get("ocr_text"))
    extraction_stages["picture_recovery"] = {
        "name": "picture_recovery",
        "chars_before": len(md),
        "chars_after": len(md),
        "char_delta": 0,
        "headings_before": len(_HEADING_RE.findall(md)),
        "headings_after": len(_HEADING_RE.findall(md)),
        "heading_delta": 0,
        "error": None,
        "regions": len(pic_results),
        "recovered": recovered_count,
    }

    return md, pic_results, extraction_stages


def _pdf_to_markdown_no_pics(pdf_path: str) -> tuple[str, list[PictureResult], dict[str, dict]]:
    """Adapter: the pymupdf4llm route recovers no picture regions and no
    per-stage provenance to match the ``(md, pics, stages)`` chain contract."""
    return pdf_to_markdown(pdf_path), [], {}


def pdf_markdown_converters() -> list[tuple[str, Callable[[str], tuple[str, list[PictureResult], dict[str, dict]]]]]:
    """Ordered ``(name, fn)`` PDF->markdown converters, per the ``PDF_CONVERTER`` env.

    Every chain callable returns ``(markdown, pic_results, extraction_stages)``.

    INDEX-01: ``pymupdf4llm`` (AGPL, fast, default) and ``docling`` (MIT,
    layout-aware, German-ligature-correct — the RFC-003 D3 / HR4 residency escape).
    The caller tries them in order and only falls back to ``page_index`` when all
    markdown converters fail. ``docling`` is listed only when importable, so a base
    install without the ``docling`` extra degrades to ``pymupdf4llm`` cleanly.

    ``docling`` is the **default** primary (it is ligature-correct on the German
    vertical and MIT-licensed, lowering AGPL exposure); set
    ``PDF_CONVERTER=pymupdf4llm`` to make the faster AGPL route primary instead, in
    which case Docling becomes the secondary markdown attempt.
    """
    import importlib.util

    from .config import ALLOW_AGPL_FALLBACK

    primary = os.getenv("PDF_CONVERTER", "docling").strip().lower()
    have_docling = importlib.util.find_spec("docling") is not None

    if not have_docling and not ALLOW_AGPL_FALLBACK:
        from .metrics import AGPL_FALLBACK_TOTAL

        AGPL_FALLBACK_TOTAL.labels(reason="blocked").inc()
        raise RuntimeError(
            "docling is not installed and ALLOW_AGPL_FALLBACK=false; "
            "either install docling (uv sync --extra docling) or set "
            "ALLOW_AGPL_FALLBACK=true"
        )

    chain: list[tuple[str, Callable[[str], tuple[str, list[PictureResult], dict[str, dict]]]]] = []
    if ALLOW_AGPL_FALLBACK:
        chain.append(("pymupdf4llm", _pdf_to_markdown_no_pics))
    if have_docling:
        if primary == "docling":
            chain.insert(0, ("docling", pdf_to_markdown_docling))
        else:
            chain.append(("docling", pdf_to_markdown_docling))
            if ALLOW_AGPL_FALLBACK:
                from .metrics import AGPL_FALLBACK_TOTAL

                AGPL_FALLBACK_TOTAL.labels(reason="operator_configured").inc()
    elif primary == "docling":
        logger.warning(
            "PDF_CONVERTER=docling but docling is not installed; install the "
            "'docling' extra (uv sync --extra docling). Falling back to pymupdf4llm."
        )
        from .metrics import AGPL_FALLBACK_TOTAL

        AGPL_FALLBACK_TOTAL.labels(reason="docling_missing").inc()
    return chain


def libreoffice_to_pdf(input_path: str) -> str:
    """Convert a DOCX/PPTX file to PDF via LibreOffice headless.

    Returns the path to the generated PDF in a temporary directory.
    The caller is responsible for cleaning up the parent directory:
        shutil.rmtree(os.path.dirname(pdf_path), ignore_errors=True)
    """
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        raise RuntimeError(
            "LibreOffice not found. Install libreoffice-headless and ensure it is on PATH."
        )
    outdir = tempfile.mkdtemp(prefix="lo_pdf_")
    # Each conversion gets its own profile dir so parallel invocations don't conflict.
    profile_dir = os.path.join(outdir, "lo_profile")
    os.makedirs(profile_dir, exist_ok=True)
    try:
        result = subprocess.run(
            [
                lo,
                f"-env:UserInstallation=file://{profile_dir}",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                outdir,
                input_path,
            ],
            capture_output=True,
            text=True,
            timeout=180,
        )
        stem = os.path.splitext(os.path.basename(input_path))[0]
        pdf_path = os.path.join(outdir, f"{stem}.pdf")
        # Check for the PDF first; a non-zero exit may be a recoverable warning
        if not os.path.isfile(pdf_path):
            pdfs = [f for f in os.listdir(outdir) if f.endswith(".pdf")]
            if pdfs:
                pdf_path = os.path.join(outdir, pdfs[0])
            elif result.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice conversion failed (exit {result.returncode}): "
                    f"{result.stderr.strip()}"
                )
            else:
                raise RuntimeError("LibreOffice did not produce a PDF file.")
        return pdf_path
    except Exception:
        shutil.rmtree(outdir, ignore_errors=True)
        raise


async def html_to_markdown_with_images(path: str, model: str) -> str:
    """Convert an HTML file to markdown, replacing <img> tags with vision-API descriptions.

    Images are described concurrently via the OpenAI vision API and inserted as
    [Image: <description>] markers at the position of the original <img> tag.
    """
    import html2text

    with open(path, encoding="utf-8", errors="replace") as f:
        html_content = f.read()

    img_pattern = re.compile(r"<img[^>]+src=[\"']([^\"']+)[\"'][^>]*/?>", re.IGNORECASE)
    srcs = img_pattern.findall(html_content)

    async def _describe(src: str) -> str:
        import openai

        from .client import get_openai_client
        from .metrics import IMAGE_DESCRIBE_FAILURES

        async def _call() -> str:
            client = get_openai_client()
            response = await client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": src}},
                            {
                                "type": "text",
                                "text": (
                                    "Describe this image concisely in 1-2 "
                                    "sentences for document context."
                                ),
                            },
                        ],
                    }
                ],
                max_tokens=150,
            )
            return response.choices[0].message.content.strip()

        try:
            return await _call()
        except (openai.RateLimitError, openai.APIConnectionError):
            # Transient failure — retry once after a short backoff.
            await asyncio.sleep(2)
            try:
                return await _call()
            except openai.APIError as retry_exc:
                logger.error(
                    "Image description failed after retry (%s): %s",
                    type(retry_exc).__name__,
                    str(retry_exc)[:200],
                )
                IMAGE_DESCRIBE_FAILURES.labels(error_type=type(retry_exc).__name__).inc()
                return "image"
            except Exception:
                # Non-OpenAI error on retry (e.g. code bug) — do not swallow it.
                raise
        except openai.APIError as exc:
            logger.error(
                "Image description failed (%s): %s",
                type(exc).__name__,
                str(exc)[:200],
            )
            IMAGE_DESCRIBE_FAILURES.labels(error_type=type(exc).__name__).inc()
            return "image"

    descriptions = await asyncio.gather(*(_describe(src) for src in srcs))

    counter = iter(range(len(descriptions)))

    def _replace(match: re.Match) -> str:
        i = next(counter, None)
        desc = descriptions[i] if i is not None else "image"
        return f"[Image: {desc}]"

    modified_html = img_pattern.sub(_replace, html_content)

    h = html2text.HTML2Text()
    h.ignore_images = True
    h.ignore_links = False
    h.body_width = 0
    return normalize_dashes(h.handle(modified_html))


def xlsx_to_markdown(path: str) -> str:
    """Convert an .xlsx workbook to markdown tables (Fix 4; openpyxl is MIT, HR4).

    Each sheet becomes a `## <sheet>` title plus one markdown table (header row +
    separator + data rows), so the existing flat-table path (route_and_extract_flat ->
    _flat_parse_table) captures every cell into row_records and Fix-2 column-stitching
    is reused. Spreadsheets carry no heading hierarchy, so the document routes flat by
    design. read_only + data_only keeps memory bounded on large books."""
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    out: list[str] = []
    try:
        for ws in wb.worksheets:
            rows = [
                ["" if c is None else str(c) for c in row] for row in ws.iter_rows(values_only=True)
            ]
            rows = [r for r in rows if any(cell.strip() for cell in r)]
            if not rows:
                continue
            width = max(len(r) for r in rows)
            out.append(f"## {ws.title}")
            header = rows[0] + [""] * (width - len(rows[0]))
            out.append("| " + " | ".join(c.replace("|", r"\|") for c in header) + " |")
            out.append("| " + " | ".join(["---"] * width) + " |")
            for r in rows[1:]:
                padded = r + [""] * (width - len(r))
                out.append("| " + " | ".join(c.replace("|", r"\|") for c in padded) + " |")
            out.append("")
    finally:
        wb.close()
    md = "\n".join(out).strip()
    if not md:
        raise RuntimeError(f"xlsx_to_markdown produced empty output for {path}")
    return md


def image_to_markdown(path: str, ocr_lang_override: list[str] | None = None) -> str:
    """OCR a scanned image (.png/.jpg/.jpeg/.tiff) to markdown (Fix 4; OCR-only, HR3).

    An image has no text layer, so full-page OCR is always on. Routes the image through
    Docling's PDF/image pipeline with the same CPU-only Tesseract options as the PDF
    path (force_full_page_ocr + Fix-5 detected language). No VLM, no LLM egress -- local
    Tesseract only (HR3). VLM stays disabled by design (RFC-004)."""
    # Docling routes InputFormat.IMAGE through the same StandardPdfPipeline as PDF, so
    # this reuses the process-lifetime _docling_converter cache (see its docstring) --
    # a fresh DocumentConverter per call leaks ~237 MB RSS (torch/models never returned).
    converter = _docling_converter(
        force_full_page_ocr=True, ocr_lang_override=ocr_lang_override, for_image=True
    )
    result = converter.convert(path)
    md = _repair_docling_tables(result.document.export_to_markdown(), doc_name=path)
    if not md or not md.strip():
        raise RuntimeError(f"image_to_markdown produced empty output for {path}")
    return normalize_dashes(md)


def rasterize_pdf_pages(pdf_path: str, dpi: int = 200) -> list[str]:
    """Rasterize each PDF page to a base64 data-URI PNG via pypdfium2 (HR4-safe)."""
    import base64
    import io

    import pypdfium2 as pdfium
    from PIL import Image

    pdoc = pdfium.PdfDocument(pdf_path)
    try:
        result: list[str] = []
        scale = dpi / 72
        for page_index in range(len(pdoc)):
            page = pdoc[page_index]
            bitmap = page.render(scale=scale)  # type: ignore[arg-type]
            pil_image: Image.Image = bitmap.to_pil()
            buf = io.BytesIO()
            pil_image.save(buf, format="PNG", optimize=True)
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            result.append(f"data:image/png;base64,{b64}")
            page.close()
        return result
    finally:
        pdoc.close()


# D4 (RFC-024): fallback rasterization backend for tesseract_ocr_pdf_pages. CMap
# corruption that crashes pypdfium2 page rendering is deterministic per-page, so a
# retry against pypdfium2 would fail identically; fitz uses a different rendering
# path (already proven for crop rasterization in _recover_picture_text) and isolates
# D7's rasterization from rasterize_pdf_pages' shared use by the VLM fallback.
_D7_FITZ_FALLBACK_ENABLED = os.getenv("D7_FITZ_FALLBACK_ENABLED", "true").strip().lower() in (
    "1",
    "true",
    "yes",
)

# RFC-029 D4 (Task 5.1) — post-export table-repair constants
# Feature flag: set to "0" to disable the repair pass entirely.
_RFC029_TABLE_DEDUP_ENABLED: bool = os.environ.get("RFC029_TABLE_DEDUP_ENABLED", "1") != "0"
# Minimum column count that must be identical before a row is collapsed.
# Rows with <= this many identical cells are left untouched to avoid collapsing
# legitimately short tables (e.g. 2-col header rows where both cols share a value).
_RFC029_TABLE_MIN_COLLAPSE_COLS: int = int(os.environ.get("RFC029_TABLE_MIN_COLLAPSE_COLS", "3"))


def rasterize_pdf_pages_fitz(pdf_path: str, dpi: int = 200) -> list[str]:
    """Rasterize each PDF page to a base64 data-URI PNG via fitz (D4, RFC-024).

    Fallback backend for ``rasterize_pdf_pages`` when pypdfium2 crashes on
    CMap-corrupt PDFs. Reuses the ``fitz.Page.get_pixmap()`` pattern already
    proven for image cropping in ``_recover_picture_text``."""
    from .config import ALLOW_AGPL_FALLBACK

    if not ALLOW_AGPL_FALLBACK:
        # RFC-034 D4 step 4: this is a literal fallback backend, so a blocked
        # invocation is a prevented AGPL fallback — meter it for observability.
        from .metrics import AGPL_FALLBACK_TOTAL

        AGPL_FALLBACK_TOTAL.labels(reason="blocked").inc()
        raise RuntimeError(
            f"cannot rasterize {pdf_path} via the fitz fallback backend: fitz "
            "(PyMuPDF, AGPL-3.0) is required and ALLOW_AGPL_FALLBACK=false"
        )
    import base64

    import fitz  # PyMuPDF, AGPL-3.0

    pdf = fitz.open(pdf_path)
    try:
        result: list[str] = []
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)
        for page in pdf:
            pix = page.get_pixmap(matrix=matrix)
            b64 = base64.b64encode(pix.tobytes("png")).decode("ascii")
            result.append(f"data:image/png;base64,{b64}")
        return result
    finally:
        pdf.close()


async def tesseract_ocr_pdf_pages(pdf_path: str, langs: list[str]) -> str:
    """Rasterize each PDF page and OCR it via local Tesseract (RFC-023 D7).

    Last-resort recovery when the VLM fallback itself crashes on a garbled
    PDF -- no LLM egress, local ``tesseract`` binary only (HR3).

    D4 (RFC-024): tries pypdfium2 (``rasterize_pdf_pages``) first; on Exception
    (e.g. CMap-corrupt PDFs that crash pypdfium2), falls back to fitz
    (``rasterize_pdf_pages_fitz``) unless disabled via
    ``D7_FITZ_FALLBACK_ENABLED=false``."""
    import base64

    try:
        page_images = await asyncio.to_thread(rasterize_pdf_pages, pdf_path)
    except Exception as exc:  # D4: fall back to fitz rasterization
        if not _D7_FITZ_FALLBACK_ENABLED:
            raise
        logger.warning(
            "rasterize_pdf_pages (pypdfium2) failed for %s (%s); falling back to fitz",
            pdf_path,
            exc,
        )
        page_images = await asyncio.to_thread(rasterize_pdf_pages_fitz, pdf_path)
    pages_text = []
    for data_uri in page_images:
        png_bytes = base64.b64decode(data_uri.split(",", 1)[1])
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as png_tmp:
            png_tmp.write(png_bytes)
            png_path = png_tmp.name
        try:
            text = await asyncio.to_thread(_tesseract_ocr_image, png_path, langs)
        finally:
            os.unlink(png_path)
        if text:
            pages_text.append(text)
    return "\n\n".join(pages_text)


async def vlm_extract_markdown(pdf_path: str, model: str | None = None) -> str:
    """Extract markdown from a PDF via a vision LLM — last-resort garble fallback."""
    import openai

    from .client import get_openai_client
    from .config import settings

    resolved_model = model or settings.vlm_model
    if resolved_model.startswith("azure/"):
        resolved_model = resolved_model[len("azure/") :]

    page_images = await asyncio.to_thread(rasterize_pdf_pages, pdf_path)
    if not page_images:
        raise RuntimeError(f"vlm_extract_markdown: no pages rasterized from {pdf_path}")

    client = get_openai_client()

    _VLM_PAGE_PROMPT = (
        "You are a document OCR assistant. Extract ALL visible text content from "
        "this scanned document page and return it as clean Markdown.\n\n"
        "Rules:\n"
        "- Preserve the document's heading hierarchy using Markdown heading levels "
        "(#, ##, ###, etc.).\n"
        "- Preserve tables as Markdown tables.\n"
        "- Preserve numbered and bulleted lists.\n"
        "- Ignore watermarks, background patterns, and page numbers.\n"
        "- If the page contains Arabic or right-to-left text, preserve the original "
        "script — do NOT transliterate.\n"
        "- Do NOT describe images; extract only text.\n"
        "- If the page is blank or contains no readable text, return exactly: "
        "<!-- blank page -->\n"
        "- Return ONLY the extracted Markdown, no commentary or wrapper."
    )

    async def _extract_page(page_idx: int, image_uri: str) -> tuple[int, str]:
        async def _call() -> str:
            response = await client.chat.completions.create(
                model=resolved_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": image_uri}},
                            {"type": "text", "text": _VLM_PAGE_PROMPT},
                        ],
                    }
                ],
                max_tokens=4096,
                temperature=0.0,
            )
            return response.choices[0].message.content.strip()

        try:
            return (page_idx, await _call())
        except (openai.RateLimitError, openai.APIConnectionError):
            await asyncio.sleep(2)
            try:
                return (page_idx, await _call())
            except Exception as retry_exc:
                logger.error("VLM page %d failed after retry: %s", page_idx + 1, retry_exc)
                return (page_idx, "")
        except Exception as exc:
            logger.error("VLM page %d extraction failed: %s", page_idx + 1, exc)
            return (page_idx, "")

    sem = asyncio.Semaphore(4)

    async def _bounded(idx: int, uri: str) -> tuple[int, str]:
        async with sem:
            return await _extract_page(idx, uri)

    results = await asyncio.gather(*[_bounded(i, u) for i, u in enumerate(page_images)])
    results_sorted = sorted(results, key=lambda r: r[0])
    page_markdowns = [md for _, md in results_sorted if md and md.strip() != "<!-- blank page -->"]

    if not page_markdowns:
        raise RuntimeError(
            f"vlm_extract_markdown: VLM returned no content for any page of {pdf_path}"
        )

    return "\n\n---\n\n".join(page_markdowns)


def docx_to_markdown(path: str) -> str:
    """Convert a DOCX file to a markdown string preserving heading hierarchy."""
    from docx import Document

    doc = Document(path)
    lines = []
    heading_map = {
        "Heading 1": "#",
        "Heading 2": "##",
        "Heading 3": "###",
        "Heading 4": "####",
        "Heading 5": "#####",
        "Heading 6": "######",
    }
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        prefix = next((v for k, v in heading_map.items() if para.style.name.startswith(k)), None)
        lines.append(f"{prefix} {text}" if prefix else text)
    return normalize_dashes("\n".join(lines))


def pptx_to_markdown(path: str) -> str:
    """Convert a PPTX file to markdown, one H1 section per slide."""
    from pptx import Presentation

    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        title = (
            title_shape.text.strip() if title_shape and title_shape.text.strip() else f"Slide {i}"
        )
        lines.append(f"# {title}")
        for shape in slide.shapes:
            if shape == title_shape or not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
        lines.append("")
    return normalize_dashes("\n".join(lines))
